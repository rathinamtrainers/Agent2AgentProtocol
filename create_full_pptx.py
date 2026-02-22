from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colors
DARK_BLUE  = RGBColor(0x1B, 0x3A, 0x6B)
ACCENT     = RGBColor(0xF0, 0xA5, 0x00)
BODY_BG    = RGBColor(0xF5, 0xF7, 0xFA)
CODE_BG    = RGBColor(0x2D, 0x2D, 0x2D)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0x4A, 0x7A, 0xB5)
LIGHT_GRAY = RGBColor(0xD0, 0xD8, 0xE8)
GREEN      = RGBColor(0x28, 0xA7, 0x45)
RED        = RGBColor(0xDC, 0x35, 0x45)
NEAR_BLACK = RGBColor(0x22, 0x22, 0x44)
PURPLE     = RGBColor(0x7B, 0x2D, 0x8B)
ORANGE     = RGBColor(0xC6, 0x6B, 0x00)
MID_BLUE   = RGBColor(0x17, 0x6A, 0xA0)
DARK_BG2   = RGBColor(0x0F, 0x22, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── helpers ────────────────────────────────────────────────────────────────
def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill, lc=None, lw=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc:
        s.line.color.rgb = lc
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h, sz=14, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return box

def multi_para(slide, lines, x, y, w, h, sz=13, color=NEAR_BLACK,
               indent_marker="    ", font="Calibri"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        if line.startswith(indent_marker):
            r.text = "       ◦  " + line.strip()
            r.font.size = Pt(sz - 1)
        else:
            r.text = "  •  " + line.strip()
            r.font.size = Pt(sz)
        r.font.color.rgb = color; r.font.name = font

def title_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.1, DARK_BLUE)
    tb(slide, title, 0.3, 0.06, 11, 0.62, sz=24, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, 0.3, 0.68, 11, 0.35, sz=12, color=ACCENT)
    rect(slide, 0, 1.1, 13.33, 0.05, ACCENT)

def section_box(slide, title, bullets, x, y, w, h, bsz=12):
    rect(slide, x, y, w, h, WHITE, LIGHT_GRAY, 1)
    rect(slide, x, y, w, 0.36, DARK_BLUE)
    tb(slide, title, x+0.1, y+0.04, w-0.2, 0.3, sz=12, bold=True, color=WHITE)
    multi_para(slide, bullets, x+0.12, y+0.42, w-0.24, h-0.5, sz=bsz)

def code_box(slide, code, x, y, w, h, sz=9.5, label=None):
    rect(slide, x, y, w, h, CODE_BG, RGBColor(0x55,0x55,0x55), 1)
    oy = y
    if label:
        rect(slide, x, y, w, 0.27, DARK_BLUE)
        tb(slide, label, x+0.1, y+0.03, w-0.2, 0.22, sz=10, color=ACCENT)
        oy = y + 0.29
        ch = h - 0.34
    else:
        oy = y + 0.08
        ch = h - 0.12
    box = slide.shapes.add_textbox(Inches(x+0.1), Inches(oy), Inches(w-0.2), Inches(ch))
    box.word_wrap = False
    tf = box.text_frame; tf.word_wrap = False
    for i, line in enumerate(code.strip().split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(sz)
        r.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        r.font.name = "Courier New"

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, DARK_BLUE)
rect(s, 0, 0, 0.5, 7.5, ACCENT)
tb(s, "Agent-to-Agent (A2A) Protocol", 1.0, 1.0, 11.8, 1.1, sz=40, bold=True, color=WHITE)
tb(s, "Complete Tutorial with Google ADK", 1.0, 2.12, 11.8, 0.65, sz=26, color=ACCENT)
rect(s, 1.0, 2.85, 10.8, 0.05, ACCENT)
tb(s, "Building a Travel Planning Multi-Agent System", 1.0, 3.0, 11.8, 0.55, sz=20, color=WHITE)
for i, pt in enumerate([
    "• Open HTTP-based standard (JSON-RPC 2.0) — works with any HTTP client",
    "• Framework-agnostic: ADK, LangChain, CrewAI, AutoGen all interoperate",
    "• Full protocol coverage: streaming, multimodal, auth, push notifications",
]):
    tb(s, pt, 1.0, 3.72 + i*0.48, 11.5, 0.45, sz=15, color=RGBColor(0xCC,0xDD,0xFF))
rect(s, 0.5, 6.72, 12.83, 0.78, DARK_BG2)
tb(s, "12 Sections  •  30+ Slides  •  Hands-on Code  •  Complete Working System",
   1.0, 6.8, 11.5, 0.45, sz=13, color=ACCENT)
notes(s, "Welcome to the A2A Protocol tutorial. We build a complete Travel Planning system with 4 specialist agents connected via the A2A protocol. All code shown is real, runnable Python.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — TABLE OF CONTENTS
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Table of Contents", "12 Sections — Full A2A Protocol Coverage")

toc = [
    ("01","What Is A2A and Why Does It Matter"),
    ("02","Agent Cards: How Agents Advertise Themselves"),
    ("03","Tasks: The Core Unit of Work"),
    ("04","Task Lifecycle: States and Transitions"),
    ("05","Streaming with Server-Sent Events (SSE)"),
    ("06","Multimodal Content: Files and Structured Data"),
    ("07","Input Required: Human-in-the-Loop"),
    ("08","Push Notifications: Fire-and-Forget Async"),
    ("09","Authentication and Security"),
    ("10","Error Handling and Resilience"),
    ("11","Multi-Agent Orchestration"),
    ("12","Running the Complete System"),
]
for i,(num,title) in enumerate(toc):
    c = 0 if i < 6 else 1
    r = i if i < 6 else i-6
    x = 0.3 if c == 0 else 6.8
    y = 1.3 + r*0.95
    rect(s, x, y, 0.58, 0.65, DARK_BLUE)
    tb(s, num, x, y+0.1, 0.58, 0.45, sz=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    rect(s, x+0.6, y, 5.82, 0.65, WHITE, LIGHT_GRAY, 1)
    tb(s, title, x+0.72, y+0.14, 5.6, 0.42, sz=12, color=NEAR_BLACK)
notes(s, "12 sections covering every aspect of A2A protocol from fundamentals to production deployment.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ARCHITECTURE DIAGRAM
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Project: Travel Planning Multi-Agent System", "Architecture Overview")

rect(s, 0.3, 1.45, 2.0, 0.7, DARK_BLUE, ACCENT, 2)
tb(s, "User / Client", 0.3, 1.55, 2.0, 0.5, sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "→", 2.35, 1.55, 0.5, 0.5, sz=26, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
rect(s, 2.88, 1.3, 3.3, 1.05, DARK_BLUE, ACCENT, 2)
tb(s, "Travel Orchestrator", 2.88, 1.37, 3.3, 0.42, sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "ADK LlmAgent + A2A Client", 2.88, 1.78, 3.3, 0.35, sz=11, color=ACCENT, align=PP_ALIGN.CENTER)

rect(s, 0.2, 2.6, 12.9, 0.22, RGBColor(0xC8,0xD5,0xEB), DARK_BLUE, 1)
tb(s, "A2A Protocol (HTTP / JSON-RPC 2.0) — all agent connections",
   0.2, 2.6, 12.9, 0.22, sz=11, color=DARK_BLUE, align=PP_ALIGN.CENTER)

agents = [
    ("Flight Agent","8001","Bearer Auth\nSync+Streaming",0.25,GREEN),
    ("Hotel Agent","8002","API Key Auth\nMultimodal PDF",3.45,MID_BLUE),
    ("Booking Agent","8003","Input-Required\nPush Notif.",6.65,PURPLE),
    ("Weather Agent","8004","SSE Streaming\nNo Auth",9.85,ORANGE),
]
for name,port,desc,x,color in agents:
    tb(s, "↓", x+0.8, 2.85, 0.7, 0.38, sz=20, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    rect(s, x, 3.28, 2.95, 1.4, color, WHITE, 1)
    tb(s, name, x, 3.32, 2.95, 0.4, sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, 3.72, 2.95, 0.22, DARK_BG2)
    tb(s, f"Port {port}", x, 3.72, 2.95, 0.22, sz=10, color=ACCENT, align=PP_ALIGN.CENTER)
    tb(s, desc, x, 3.96, 2.95, 0.65, sz=11, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, 0.2, 5.02, 12.9, 0.6, WHITE, LIGHT_GRAY, 1)
tb(s, "Tech Stack:  HTTP/HTTPS  •  JSON-RPC 2.0  •  Server-Sent Events  •  Google ADK  •  FastAPI  •  Docker Compose",
   0.4, 5.1, 12.5, 0.42, sz=12, color=NEAR_BLACK)
notes(s, "Architecture overview. The Travel Orchestrator is an ADK LlmAgent that uses A2AClient as a tool. It discovers capabilities from each agent's Agent Card, fans out parallel requests, and aggregates results.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SECTION 1: WHAT IS A2A
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 1: What Is A2A and Why Does It Matter", "The open standard for agent interoperability")

rect(s, 0.3, 1.25, 5.9, 2.8, WHITE, LIGHT_GRAY, 1)
rect(s, 0.3, 1.25, 5.9, 0.38, RED)
tb(s, "Before A2A — The Problem", 0.4, 1.27, 5.7, 0.32, sz=13, bold=True, color=WHITE)
multi_para(s, [
    "Every multi-agent system needed custom glue code",
    "Framework lock-in: ADK agents couldn't talk to LangChain agents",
    "No standard for discovery, task lifecycle, or streaming",
    "Bespoke auth, error codes, and message formats per project",
    "Hard to test, hard to swap out individual agents",
    "No ecosystem — every team reinvented the wheel",
], 0.42, 1.7, 5.68, 2.3, sz=12, color=NEAR_BLACK)

rect(s, 6.5, 1.25, 6.53, 2.8, WHITE, LIGHT_GRAY, 1)
rect(s, 6.5, 1.25, 6.53, 0.38, GREEN)
tb(s, "After A2A — The Solution", 6.6, 1.27, 6.35, 0.32, sz=13, bold=True, color=WHITE)
multi_para(s, [
    "Open HTTP-based protocol — any HTTP client works",
    "JSON-RPC 2.0 as the message format",
    "Any compliant agent talks to any other agent",
    "Standard task lifecycle, streaming, and auth",
    "Machine-readable Agent Cards for auto-discovery",
    "Built on proven web standards (HTTP, SSE, JSON)",
], 6.62, 1.7, 6.35, 2.3, sz=12, color=NEAR_BLACK)

# Key components row
rect(s, 0.3, 4.3, 12.73, 0.38, DARK_BLUE)
tb(s, "Key A2A Components", 0.5, 4.33, 12.0, 0.3, sz=13, bold=True, color=WHITE)
components = [
    ("Agent Cards","Machine-readable\nagent metadata"),
    ("Tasks","Unit of work\nwith lifecycle"),
    ("Messages","Structured\ncommunication"),
    ("Parts","Text, File,\nor Data content"),
    ("Artifacts","Task output\ndeliverables"),
]
for i,(name,desc) in enumerate(components):
    x = 0.3 + i*2.55
    rect(s, x, 4.7, 2.42, 1.3, LIGHT_BLUE, DARK_BLUE, 1)
    tb(s, name, x, 4.73, 2.42, 0.4, sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, desc, x, 5.15, 2.42, 0.75, sz=11, color=WHITE, align=PP_ALIGN.CENTER)
notes(s, "A2A solves the agent interoperability problem by defining a standard protocol. Before A2A, every system needed custom integration code. After A2A, any compliant agent can communicate with any other regardless of the underlying framework.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — A2A KEY CONCEPTS DEEP DIVE
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 1: A2A Core Architecture", "HTTP + JSON-RPC 2.0 + SSE")

rect(s, 0.3, 1.25, 8.2, 4.7, WHITE, LIGHT_GRAY, 1)
rect(s, 0.3, 1.25, 8.2, 0.38, DARK_BLUE)
tb(s, "Wire Protocol — What Actually Goes Over the Network", 0.4, 1.27, 8.0, 0.32, sz=12, bold=True, color=WHITE)
code_box(s, """\
# Standard JSON-RPC 2.0 request
POST /  HTTP/1.1
Content-Type: application/json

{
  "jsonrpc": "2.0",
  "id": "req-001",
  "method": "tasks/send",
  "params": {
    "id": "task-uuid-1234",
    "message": {
      "role": "user",
      "parts": [{"type": "text", "text": "Find flights to Paris"}]
    }
  }
}

# JSON-RPC 2.0 success response
{"jsonrpc":"2.0","id":"req-001","result": {...}}

# JSON-RPC 2.0 error response
{"jsonrpc":"2.0","id":"req-001","error":{"code":-32602,"message":"Invalid params"}}
""", 0.35, 1.68, 8.1, 4.22, sz=9.5)

section_box(s, "A2A Endpoints", [
    "POST /  →  JSON-RPC method dispatch",
    "GET /.well-known/agent.json  →  Agent Card",
    "tasks/send  →  send a task synchronously",
    "tasks/sendSubscribe  →  send + stream SSE",
    "tasks/get  →  poll task status",
    "tasks/cancel  →  cancel a running task",
    "tasks/pushNotification/set  →  register webhook",
    "tasks/pushNotification/get  →  get webhook config",
], 8.65, 1.25, 4.45, 4.7, bsz=11)
notes(s, "A2A uses standard JSON-RPC 2.0 over HTTP POST to /. The agent.json endpoint serves the Agent Card for discovery. All task operations use the tasks/* method namespace. SSE streaming is triggered by tasks/sendSubscribe.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SECTION 2: AGENT CARDS
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 2: Agent Cards", "Machine-readable agent discovery at /.well-known/agent.json")

code_box(s, """\
# Flight Agent Card  — /.well-known/agent.json
{
  "name": "Flight Search Agent",
  "description": "Searches and books flights globally",
  "url": "http://localhost:8001/",
  "version": "1.0.0",
  "capabilities": {
    "streaming": true,
    "pushNotifications": false,
    "stateTransitionHistory": true
  },
  "authentication": {
    "schemes": ["Bearer"]
  },
  "defaultInputModes":  ["text/plain"],
  "defaultOutputModes": ["application/json"],
  "skills": [
    {
      "id": "search_flights",
      "name": "Search Flights",
      "description": "Find available flights between cities",
      "inputModes":  ["text/plain"],
      "outputModes": ["application/json"],
      "tags": ["travel","flights","search"],
      "examples": ["Find flights from NYC to Paris on Dec 15"]
    }
  ]
}
""", 0.3, 1.22, 6.55, 6.0, sz=9.5, label="Flight Agent Card (/.well-known/agent.json)")

section_box(s, "Agent Card Fields", [
    "name — human-readable agent name",
    "description — what this agent does",
    "url — base URL for all A2A requests",
    "version — semantic version string",
    "capabilities — feature flags",
    "authentication — required auth schemes",
    "defaultInputModes — accepted MIME types",
    "defaultOutputModes — produced MIME types",
    "skills[] — individual capabilities",
], 7.1, 1.22, 5.9, 3.4, bsz=11)

section_box(s, "Capabilities Flags", [
    "streaming: true → supports tasks/sendSubscribe",
    "pushNotifications: true → supports webhooks",
    "stateTransitionHistory: true → full audit trail",
], 7.1, 4.72, 5.9, 1.55, bsz=11)

section_box(s, "Auth Schemes", [
    "Bearer — HTTP Authorization: Bearer <token>",
    "ApiKey — custom header (e.g. X-API-Key)",
    "OAuth2 — standard OAuth2 flows",
    "None   — no auth required",
], 7.1, 6.35, 5.9, 1.8, bsz=11)
notes(s, "Every A2A agent MUST serve its Agent Card at /.well-known/agent.json. Clients fetch this before sending any requests to discover capabilities, required auth, input/output modes, and available skills. This enables auto-discovery in orchestrators.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AGENT CARD: HOTEL AGENT EXAMPLE
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 2: Agent Cards — Hotel & Booking Examples", "Showing different auth schemes and capabilities")

code_box(s, """\
# Hotel Agent Card
{
  "name": "Hotel Booking Agent",
  "url": "http://localhost:8002/",
  "capabilities": {
    "streaming": false,
    "pushNotifications": false
  },
  "authentication": {
    "schemes": ["ApiKey"],
    "credentials": {"header": "X-API-Key"}
  },
  "skills": [{
    "id": "search_hotels",
    "name": "Search Hotels",
    "inputModes":  ["text/plain","application/pdf"],
    "outputModes": ["application/json"],
    "description": "Search hotels; accepts PDF brochures"
  }]
}
""", 0.3, 1.22, 6.3, 4.5, sz=9.5, label="Hotel Agent — API Key + PDF input")

code_box(s, """\
# Booking Agent Card
{
  "name": "Booking Confirmation Agent",
  "url": "http://localhost:8003/",
  "capabilities": {
    "streaming": false,
    "pushNotifications": true
  },
  "authentication": {"schemes": []},
  "skills": [{
    "id": "confirm_booking",
    "name": "Confirm Booking",
    "description": "Confirms travel bookings; may require "
                   "additional user input (seat, meal)",
    "inputModes":  ["application/json"],
    "outputModes": ["application/json"]
  }]
}
""", 0.3, 5.82, 6.3, 1.8, sz=9.5, label="Booking Agent — push notifications")

section_box(s, "Skills Array — Key Fields", [
    "id          — unique identifier for routing",
    "name        — human-readable skill name",
    "description — what the skill does (LLM-readable)",
    "inputModes  — list of accepted MIME types",
    "outputModes — list of produced MIME types",
    "tags        — for skill discovery/filtering",
    "examples    — example prompts for LLM routing",
], 6.8, 1.22, 6.2, 3.2, bsz=11)

section_box(s, "How Orchestrators Use Agent Cards", [
    "1. Fetch all agent cards at startup",
    "2. Match user intent to skill descriptions",
    "3. Check capability flags (streaming? push?)",
    "4. Read auth scheme and inject credentials",
    "5. Use inputModes to format the request",
    "6. Route subtask to the best-matching agent",
], 6.8, 4.52, 6.2, 3.2, bsz=11)
notes(s, "Different agents declare different capabilities. Hotel Agent uses API Key auth and accepts PDF input. Booking Agent has no auth but supports push notifications. Orchestrators read these cards to know how to call each agent correctly.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SECTION 3: TASKS
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 3: Tasks — The Core Unit of Work", "Every interaction is a Task sent via tasks/send")

code_box(s, """\
# TaskSendParams — what you send
{
  "id":        "550e8400-e29b-41d4-a716-446655440000",  # client-generated UUID
  "sessionId": "session-travel-2024-001",               # groups related tasks
  "message": {
    "role": "user",                                       # "user" | "agent"
    "parts": [
      {"type": "text", "text": "Find me flights to Paris on Dec 15"}
    ]
  }
}
""", 0.3, 1.22, 6.4, 2.85, sz=9.5, label="Sending a Task — POST tasks/send")

code_box(s, """\
# Task Response
{
  "id":        "550e8400-e29b-41d4-a716-446655440000",
  "sessionId": "session-travel-2024-001",
  "status": {
    "state":     "completed",
    "message":   null,
    "timestamp": "2024-12-15T10:30:00Z"
  },
  "artifacts": [{
    "name":        "flight_results",
    "description": "Available flights",
    "index":       0,
    "lastChunk":   true,
    "parts": [{"type": "data", "data": {"flights": [...]}}]
  }],
  "history": [...]
}
""", 0.3, 4.17, 6.4, 3.5, sz=9.5, label="Task Response")

section_box(s, "Part Types", [
    "TextPart  — plain text content",
    "    type: 'text'",
    "    text: str",
    "FilePart  — binary or hosted files",
    "    type: 'file'",
    "    mimeType: str",
    "    bytes: base64  OR  uri: str",
    "DataPart  — structured JSON",
    "    type: 'data'",
    "    data: dict",
], 6.95, 1.22, 6.05, 4.0, bsz=11)

section_box(s, "Artifact Fields", [
    "name        — artifact identifier",
    "description — human-readable description",
    "parts[]     — content (same Part types)",
    "index       — chunk sequence number",
    "lastChunk   — true when stream complete",
    "append      — true if chunk extends prior",
], 6.95, 5.32, 6.05, 2.35, bsz=11)
notes(s, "Every A2A interaction is a Task. The client generates a UUID task ID. SessionId groups multiple tasks in a conversation. Parts carry the actual content as text, files (inline bytes or hosted URI), or structured data. Artifacts are the outputs from the agent.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SECTION 3: TASKS — PYTHON SDK
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 3: Tasks — Python A2A Client", "Using the google-a2a-sdk to send tasks")

code_box(s, """\
import uuid
from a2a.client import A2AClient
from a2a.types  import (
    SendTaskRequest, TaskSendParams,
    Message, TextPart
)

# Initialize client (reads Agent Card automatically)
client = A2AClient(url="http://localhost:8001/")

# Build and send a task
task_id = str(uuid.uuid4())
request = SendTaskRequest(
    id=task_id,
    params=TaskSendParams(
        id=task_id,
        sessionId="session-001",
        message=Message(
            role="user",
            parts=[TextPart(type="text",
                            text="Find flights NYC→Paris, Dec 15")]
        )
    )
)

response = await client.send_task(request)

# Access results
task   = response.result
status = task.status.state        # "completed"
arts   = task.artifacts           # list of Artifact
text   = arts[0].parts[0].text    # response text
data   = arts[0].parts[0].data    # structured dict
""", 0.3, 1.22, 7.0, 5.9, sz=9.5, label="tasks/send — Python SDK")

section_box(s, "TaskSendParams Fields", [
    "id        — client-generated UUID (required)",
    "message   — the Message object (required)",
    "sessionId — optional session grouping",
    "historyLength — how many history turns to return",
    "metadata  — arbitrary client metadata dict",
    "pushNotification — webhook config (Section 8)",
], 7.55, 1.22, 5.45, 3.2, bsz=11)

section_box(s, "Client Configuration", [
    "url      — agent base URL",
    "headers  — default request headers",
    "timeout  — request timeout seconds",
    "Reads Agent Card to validate capabilities",
    "Handles JSON-RPC envelope automatically",
    "Raises A2AClientHTTPError on 4xx/5xx",
], 7.55, 4.52, 5.45, 2.75, bsz=11)
notes(s, "The A2A Python SDK provides an async client. You create a Message with Parts, wrap it in TaskSendParams, and call send_task(). The SDK handles JSON-RPC framing. The client can also auto-fetch the Agent Card to validate you are sending the right content types.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SECTION 4: TASK LIFECYCLE
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 4: Task Lifecycle — States and Transitions", "From submitted to completed/failed/canceled/input-required")

# State diagram using shapes
states = [
    ("submitted", 1.0, 2.2, LIGHT_BLUE, "Initial state\nright after tasks/send"),
    ("working",   4.2, 2.2, RGBColor(0x17,0x80,0x60), "Agent is\nprocessing"),
    ("completed", 7.4, 2.2, GREEN, "Normal\ntermination"),
    ("failed",    7.4, 3.8, RED, "Error\noccurred"),
    ("canceled",  7.4, 5.4, ORANGE, "Client called\ntasks/cancel"),
    ("input-required", 4.2, 5.0, PURPLE, "Agent needs\nmore info"),
]
for name, x, y, color, desc in states:
    rect(s, x, y, 2.2, 0.8, color, WHITE, 2)
    tb(s, name, x, y+0.08, 2.2, 0.38, sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, desc, x, y+0.45, 2.2, 0.4, sz=9, color=WHITE, align=PP_ALIGN.CENTER)

# Arrows (text-based)
arrows = [
    (3.2, 2.47, "→"),
    (6.4, 2.47, "→"),
    (6.4, 4.07, "→"),
    (6.4, 5.67, "→"),
    (5.55, 3.95, "↓"),
    (5.55, 4.75, "↑ resume"),
]
for x, y, sym in arrows:
    tb(s, sym, x, y, 0.8, 0.45, sz=18, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# Endpoint table
rect(s, 0.3, 6.12, 12.7, 0.62, WHITE, LIGHT_GRAY, 1)
endpoints = [
    ("tasks/send","Submit new task"),
    ("tasks/get","Poll current status"),
    ("tasks/cancel","Request cancellation"),
    ("tasks/sendSubscribe","Submit + stream updates"),
]
tb(s, "Endpoints:", 0.45, 6.18, 1.8, 0.42, sz=11, bold=True, color=DARK_BLUE)
for i,(ep,desc) in enumerate(endpoints):
    rect(s, 2.3+i*2.65, 6.18, 2.5, 0.42, DARK_BLUE)
    tb(s, f"{ep}\n{desc}", 2.3+i*2.65, 6.18, 2.5, 0.42, sz=9, color=WHITE, align=PP_ALIGN.CENTER)
notes(s, "Task states: submitted (immediately after tasks/send), working (agent processing), completed (success), failed (error), canceled (client called tasks/cancel), input-required (agent needs more info from user). The client can poll with tasks/get or use tasks/sendSubscribe for streaming updates.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TASK LIFECYCLE: POLLING AND STATE CODE
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 4: Task Lifecycle — Polling and State Handling", "Polling with tasks/get and state machine patterns")

code_box(s, """\
import asyncio
from a2a.client import A2AClient
from a2a.types  import GetTaskRequest, TaskQueryParams

async def poll_until_done(client: A2AClient, task_id: str,
                          poll_interval: float = 2.0):
    \"\"\"Poll tasks/get until task reaches a terminal state.\"\"\"
    terminal = {"completed", "failed", "canceled"}

    while True:
        resp = await client.get_task(
            GetTaskRequest(
                id=task_id,
                params=TaskQueryParams(id=task_id, historyLength=0)
            )
        )
        task  = resp.result
        state = task.status.state

        print(f"Task {task_id}: state={state}")

        if state in terminal:
            return task

        if state == "input-required":
            # Handle in Section 7
            return task

        await asyncio.sleep(poll_interval)

# Usage
task  = await poll_until_done(client, task_id)
state = task.status.state

if state == "completed":
    process_artifacts(task.artifacts)
elif state == "failed":
    print(f"Failed: {task.status.message}")
elif state == "canceled":
    print("Task was canceled")
""", 0.3, 1.22, 7.1, 5.95, sz=9.5, label="Polling Pattern — tasks/get")

section_box(s, "TaskStatus Object", [
    "state     — current state string",
    "message   — optional Message from agent",
    "timestamp — ISO-8601 UTC timestamp",
    "",
    "Terminal states (no more transitions):",
    "    completed  — success, artifacts ready",
    "    failed     — error, check status.message",
    "    canceled   — client-requested cancel",
    "",
    "Non-terminal (keep polling):",
    "    submitted  — queued, not started yet",
    "    working    — actively processing",
    "    input-required — waiting for user",
], 7.55, 1.22, 5.45, 5.95, bsz=11)
notes(s, "For synchronous tasks use tasks/get to poll. Check status.state and loop until you hit a terminal state. Always handle input-required separately (Section 7). For long-running tasks consider push notifications (Section 8) instead of tight polling.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SECTION 5: STREAMING SSE
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 5: Streaming with Server-Sent Events (SSE)", "Real-time updates via tasks/sendSubscribe")

code_box(s, """\
from a2a.client import A2AClient
from a2a.types  import (
    SendTaskStreamingRequest, TaskSendParams,
    Message, TextPart
)

async def stream_flight_search(client: A2AClient):
    \"\"\"Stream results from Flight Agent using SSE.\"\"\"
    request = SendTaskStreamingRequest(
        id="stream-task-001",
        params=TaskSendParams(
            id="stream-task-001",
            message=Message(
                role="user",
                parts=[TextPart(type="text",
                                text="Find all flights NYC→Paris Dec 15")]
            )
        )
    )

    # tasks/sendSubscribe returns an async generator
    async for event in client.send_task_streaming(request):
        if event.result is None:
            continue

        result = event.result

        # TaskStatusUpdateEvent
        if hasattr(result, 'status'):
            state = result.status.state
            print(f"Status update: {state}")
            if state in ("completed","failed","canceled"):
                break

        # TaskArtifactUpdateEvent
        elif hasattr(result, 'artifact'):
            art = result.artifact
            for part in art.parts:
                if hasattr(part, 'text'):
                    print(part.text, end='', flush=True)
            if art.lastChunk:
                print()   # final newline
                break
""", 0.3, 1.22, 7.35, 6.0, sz=9.5, label="tasks/sendSubscribe — SSE Streaming")

section_box(s, "SSE Event Types", [
    "TaskStatusUpdateEvent",
    "    status.state  — new task state",
    "    status.message — optional agent message",
    "    final: true when stream ends",
    "",
    "TaskArtifactUpdateEvent",
    "    artifact.parts — content chunk",
    "    artifact.index — chunk sequence",
    "    artifact.lastChunk — stream done",
    "    artifact.append — extends prior chunk",
], 7.65, 1.22, 5.35, 4.5, bsz=11)

section_box(s, "Agent Card Requirement", [
    "capabilities.streaming: true",
    "Clients check this BEFORE subscribing",
    "",
    "Used By:",
    "    Flight Agent  — search results",
    "    Weather Agent — live forecast stream",
], 7.65, 5.82, 5.35, 1.9, bsz=11)
notes(s, "SSE streaming uses tasks/sendSubscribe. The server sends a stream of Server-Sent Events. Two event types: TaskStatusUpdateEvent (state changes) and TaskArtifactUpdateEvent (content chunks). The lastChunk flag on artifacts signals end of stream. Agent must advertise streaming:true in its Agent Card.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SSE: SERVER IMPLEMENTATION
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 5: Streaming — Server-Side Implementation", "How the Flight Agent generates SSE events")

code_box(s, """\
from a2a.server.agent_execution import AgentExecutor
from a2a.server.events import EventQueue
from a2a.types import (
    TaskStatusUpdateEvent, TaskArtifactUpdateEvent,
    TaskStatus, Artifact, TextPart
)

class FlightAgentExecutor(AgentExecutor):

    async def execute(self, context, event_queue: EventQueue):
        query = context.get_user_input()

        # Emit working status
        await event_queue.enqueue_event(
            TaskStatusUpdateEvent(
                status=TaskStatus(state="working"),
                final=False
            )
        )

        # Stream results chunk by chunk
        flights = await self._search_flights(query)
        for i, flight in enumerate(flights):
            is_last = (i == len(flights) - 1)
            await event_queue.enqueue_event(
                TaskArtifactUpdateEvent(
                    artifact=Artifact(
                        name="flights",
                        index=i,
                        lastChunk=is_last,
                        parts=[TextPart(
                            type="text",
                            text=f"Flight {i+1}: {flight}"
                        )]
                    )
                )
            )

        # Emit completed status
        await event_queue.enqueue_event(
            TaskStatusUpdateEvent(
                status=TaskStatus(state="completed"),
                final=True
            )
        )
""", 0.3, 1.22, 7.55, 6.0, sz=9.5, label="FlightAgentExecutor — Server-side SSE")

section_box(s, "EventQueue Pattern", [
    "AgentExecutor.execute() receives EventQueue",
    "Call event_queue.enqueue_event() to stream",
    "Status events: state transitions",
    "Artifact events: content chunks",
    "final=True signals end of SSE stream",
    "",
    "Framework handles:",
    "    SSE formatting (data: ... \\n\\n)",
    "    Content-Type: text/event-stream",
    "    HTTP chunked transfer encoding",
    "    Client disconnect detection",
], 8.0, 1.22, 5.0, 6.0, bsz=11)
notes(s, "On the server side, your AgentExecutor receives an EventQueue. Call enqueue_event() with either TaskStatusUpdateEvent or TaskArtifactUpdateEvent. The A2A framework handles all SSE formatting, HTTP headers, and streaming. Set final=True on the last status event to close the stream.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — SECTION 6: MULTIMODAL
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 6: Multimodal Content — Files and Structured Data", "FilePart (bytes/URI) and DataPart (structured JSON)")

code_box(s, """\
import base64
from a2a.types import (
    Message, FilePart, FileContent, DataPart, TextPart
)

# ── FilePart: inline bytes ─────────────────────────────
with open("hotel_brochure.pdf", "rb") as f:
    pdf_bytes = base64.b64encode(f.read()).decode()

pdf_part = FilePart(
    type="file",
    file=FileContent(
        mimeType="application/pdf",
        bytes=pdf_bytes       # base64-encoded
    )
)

# ── FilePart: hosted URI ───────────────────────────────
uri_part = FilePart(
    type="file",
    file=FileContent(
        mimeType="image/jpeg",
        uri="https://cdn.example.com/hotel_lobby.jpg"
    )
)

# ── DataPart: structured JSON ──────────────────────────
data_part = DataPart(
    type="data",
    data={
        "checkin":  "2024-12-15",
        "checkout": "2024-12-20",
        "guests":   2,
        "preferences": ["ocean view", "king bed"]
    }
)

# ── Send multimodal message ────────────────────────────
message = Message(
    role="user",
    parts=[
        TextPart(type="text", text="Search based on this brochure"),
        pdf_part,
        data_part,
    ]
)
""", 0.3, 1.22, 7.45, 6.0, sz=9.5, label="Multimodal Parts — Python")

section_box(s, "FilePart Variants", [
    "Inline bytes (small files < 1MB):",
    "    bytes  — base64-encoded content",
    "    mimeType — MIME type string",
    "",
    "Hosted URI (large files):",
    "    uri    — publicly accessible URL",
    "    mimeType — MIME type string",
    "",
    "Common MIME types:",
    "    application/pdf",
    "    image/jpeg, image/png",
    "    audio/wav, video/mp4",
], 7.9, 1.22, 5.1, 4.5, bsz=11)

section_box(s, "Hotel Agent Skill Declaration", [
    "inputModes: [",
    '    "text/plain",',
    '    "application/pdf"',
    "]",
    "Declared in Agent Card skill",
    "Orchestrator reads & formats accordingly",
    "Returns DataPart with amenities dict",
], 7.9, 5.82, 5.1, 1.9, bsz=11)
notes(s, "A2A supports three Part types: TextPart for plain text, FilePart for binary content (inline base64 bytes or hosted URI), and DataPart for structured JSON. The Hotel Agent advertises application/pdf as an accepted input mode in its Agent Card skill definition, so the orchestrator knows it can send PDF files.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — SECTION 7: INPUT REQUIRED
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 7: Input Required — Human-in-the-Loop", "Agent pauses and requests additional information")

code_box(s, """\
# ── STEP 1: Send initial booking request ──────────────
task_id = str(uuid.uuid4())
resp = await client.send_task(SendTaskRequest(
    id=task_id,
    params=TaskSendParams(
        id=task_id,
        sessionId="session-001",
        message=Message(role="user", parts=[
            DataPart(type="data", data=flight_data)
        ])
    )
))
task = resp.result

# ── STEP 2: Handle input-required loop ────────────────
while task.status.state == "input-required":
    # Agent's question is in status.message
    question = task.status.message.parts[0].text
    print(f"Agent asks: {question}")
    # e.g. "Which seat class? (economy/business/first)"

    # Get answer from user (or LLM)
    user_answer = input("Your answer: ")

    # RESUME the SAME task by reusing task_id
    resp = await client.send_task(SendTaskRequest(
        id=task_id,                    # <── same task ID!
        params=TaskSendParams(
            id=task_id,
            sessionId="session-001",   # <── same session
            message=Message(role="user", parts=[
                TextPart(type="text", text=user_answer)
            ])
        )
    ))
    task = resp.result

# ── STEP 3: Process final result ──────────────────────
if task.status.state == "completed":
    booking_ref = task.artifacts[0].parts[0].data
    print(f"Booking confirmed: {booking_ref['reference']}")
""", 0.3, 1.22, 7.45, 6.05, sz=9.0, label="Booking Agent — input-required loop")

section_box(s, "Booking Agent Multi-Step Flow", [
    "Step 1 — Send flight/hotel selection",
    "    → state: input-required",
    '    → "Which seat class?"',
    "Step 2 — Send seat class (economy)",
    "    → state: input-required",
    '    → "Meal preference?"',
    "Step 3 — Send meal preference",
    "    → state: input-required",
    '    → "Loyalty number? (skip=none)"',
    "Step 4 — Send loyalty number",
    "    → state: completed",
    "    → booking_reference artifact",
], 7.9, 1.22, 5.1, 6.05, bsz=11)
notes(s, "When an agent needs more information it returns state=input-required. The agent's question is in status.message. The client MUST reuse the same task_id when resuming — this continues the same task. SessionId groups the conversation. This pattern enables human-in-the-loop or multi-step clarification flows.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — SECTION 8: PUSH NOTIFICATIONS
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 8: Push Notifications — Fire-and-Forget Async", "Agent POSTs to your webhook when task completes")

code_box(s, """\
from a2a.types import (
    SetTaskPushNotificationRequest,
    TaskPushNotificationConfig,
    PushNotificationConfig,
    PushNotificationAuthenticationInfo
)

# ── Register webhook BEFORE sending the task ──────────
await client.set_task_push_notification(
    SetTaskPushNotificationRequest(
        id=task_id,
        params=TaskPushNotificationConfig(
            id=task_id,
            pushNotificationConfig=PushNotificationConfig(
                url="https://myserver.com/webhook/a2a",
                token="secret-verify-token-abc123",
                authentication=PushNotificationAuthenticationInfo(
                    schemes=["Bearer"],
                    credentials="my-webhook-bearer-token"
                )
            )
        )
    )
)

# ── Send the long-running task ─────────────────────────
await client.send_task(SendTaskRequest(
    id=task_id,
    params=TaskSendParams(id=task_id, message=...)
))
print("Task submitted — webhook will fire when done")
# Client is free to do other work now!
""", 0.3, 1.22, 7.1, 5.6, sz=9.0, label="Register Webhook + Send Task")

code_box(s, """\
# webhook_receiver.py — FastAPI server
from fastapi import FastAPI, Request, HTTPException
import hmac, hashlib

app = FastAPI()

@app.post("/webhook/a2a")
async def receive_push(request: Request):
    # Validate token header
    token = request.headers.get("X-A2A-Notification-Token")
    if not hmac.compare_digest(token or "", EXPECTED_TOKEN):
        raise HTTPException(status_code=401)

    payload = await request.json()
    task_id = payload["id"]
    state   = payload["status"]["state"]
    print(f"Push notification: task={task_id} state={state}")

    if state == "completed":
        artifacts = payload.get("artifacts", [])
        await process_result(task_id, artifacts)
    return {"ok": True}
""", 0.3, 6.92, 7.1, 1.75, sz=8.5, label="Webhook Receiver (FastAPI)")

section_box(s, "PushNotificationConfig Fields", [
    "url    — your HTTPS webhook endpoint",
    "token  — verification token (you set it)",
    "authentication — auth for agent to call you",
    "",
    "Agent Card requirement:",
    "    capabilities.pushNotifications: true",
    "",
    "Agent POSTs full Task JSON to webhook",
    "Validate token to prevent spoofing",
    "Return 200 OK to acknowledge receipt",
    "",
    "tasks/pushNotification/get — retrieve config",
], 7.55, 1.22, 5.45, 5.05, bsz=11)

section_box(s, "When to Use Push Notifications", [
    "Tasks taking minutes or hours",
    "Batch processing pipelines",
    "When client cannot maintain connection",
    "Serverless / event-driven architectures",
    "Replace tight polling loops",
], 7.55, 6.35, 5.45, 1.35, bsz=11)
notes(s, "Push notifications are for long-running tasks where polling is impractical. The client registers a webhook URL via tasks/pushNotification/set BEFORE sending the task. The agent POSTs the full task JSON to the webhook when the task completes. Always validate the token header to prevent spoofed notifications.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — SECTION 9: AUTHENTICATION
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 9: Authentication and Security", "Agent Cards declare requirements; clients inject credentials")

code_box(s, """\
import httpx
from a2a.client import A2AClient, A2ACardResolver

# ── Pattern 1: Bearer Token (Flight Agent) ────────────
resolver = A2ACardResolver(
    httpx_client=httpx.AsyncClient(),
    base_url="http://localhost:8001/"
)
card = await resolver.get_agent_card()
# card.authentication.schemes == ["Bearer"]

flight_client = A2AClient(
    httpx_client=httpx.AsyncClient(headers={
        "Authorization": f"Bearer {FLIGHT_API_TOKEN}"
    }),
    agent_card=card
)

# ── Pattern 2: API Key (Hotel Agent) ──────────────────
hotel_card = await A2ACardResolver(
    httpx_client=httpx.AsyncClient(),
    base_url="http://localhost:8002/"
).get_agent_card()
# hotel_card.authentication.credentials.header == "X-API-Key"

hotel_client = A2AClient(
    httpx_client=httpx.AsyncClient(headers={
        "X-API-Key": HOTEL_API_KEY
    }),
    agent_card=hotel_card
)

# ── Handle auth errors ────────────────────────────────
from a2a.client.errors import A2AClientHTTPError

try:
    response = await flight_client.send_task(request)
except A2AClientHTTPError as e:
    if e.status_code == 401:
        print("Token expired — refresh and retry")
    elif e.status_code == 403:
        print("Insufficient permissions")
    else:
        raise
""", 0.3, 1.22, 7.35, 6.0, sz=9.0, label="Client-Side Authentication — Python")

section_box(s, "Auth Scheme Summary", [
    "Bearer Token (Flight Agent 8001):",
    "    Authorization: Bearer <token>",
    "    Server validates via middleware",
    "    401 if missing/invalid",
    "",
    "API Key (Hotel Agent 8002):",
    "    X-API-Key: <key>",
    "    Header name in Agent Card credentials",
    "    403 if key invalid/missing",
    "",
    "No Auth (Booking/Weather):",
    "    No credentials required",
    "    Still validate input via schema",
], 7.8, 1.22, 5.2, 4.5, bsz=11)

section_box(s, "Server-Side Middleware", [
    "Reads Agent Card authentication schemes",
    "Validates on EVERY incoming request",
    "Returns JSON-RPC error on failure:",
    '    HTTP 401 for missing credentials',
    '    HTTP 403 for invalid credentials',
    "Never skip validation in production",
], 7.8, 5.82, 5.2, 1.9, bsz=11)
notes(s, "Security pattern: clients fetch the Agent Card to discover required auth scheme, then inject the right credentials into every request. Use Bearer tokens for user-scoped auth and API Keys for service-to-service. Always handle A2AClientHTTPError and retry on 401 with refreshed tokens.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — SECTION 10: ERROR HANDLING
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 10: Error Handling and Resilience", "JSON-RPC error codes and recovery patterns")

# Error codes table
rect(s, 0.3, 1.22, 12.7, 0.38, DARK_BLUE)
for i,(w,label) in enumerate([(1.2,"Code"),(3.0,"Name"),(7.5,"Description")]):
    tb(s, label, 0.4+sum([1.2,3.0,7.5][:i] if i else []), 1.24, w, 0.3,
       sz=11, bold=True, color=WHITE)

error_rows = [
    ("-32700","Parse Error","Invalid JSON in request body"),
    ("-32600","Invalid Request","JSON-RPC envelope malformed"),
    ("-32601","Method Not Found","e.g. 'tasks/unknown' does not exist"),
    ("-32602","Invalid Params","Required field missing or wrong type"),
    ("-32603","Internal Error","Unhandled server exception"),
    ("-32001","TaskNotFound","Task ID does not exist on server"),
    ("-32002","TaskNotCancelable","Task already in terminal state"),
    ("-32003","PushNotif.NotSupported","Agent Card: pushNotifications:false"),
    ("-32004","UnsupportedOperation","Method exists but agent disables it"),
    ("-32005","ContentTypeNotSupported","MIME type not in Agent Card"),
]
for i,(code,name,desc) in enumerate(error_rows):
    bg = WHITE if i%2==0 else RGBColor(0xEE,0xF2,0xF8)
    rect(s, 0.3, 1.62+i*0.43, 12.7, 0.43, bg, LIGHT_GRAY, 1)
    tb(s, code, 0.42, 1.65+i*0.43, 1.1, 0.35, sz=11, color=RED, font="Courier New")
    tb(s, name, 1.55, 1.65+i*0.43, 2.9, 0.35, sz=11, bold=True, color=NEAR_BLACK)
    tb(s, desc, 4.48, 1.65+i*0.43, 8.4, 0.35, sz=11, color=NEAR_BLACK)
notes(s, "A2A uses standard JSON-RPC 2.0 error codes (-32700 to -32603) plus A2A-specific codes starting at -32001. Always check the error code in catch blocks to decide whether to retry, surface to user, or gracefully degrade.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — ERROR HANDLING: PATTERNS
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 10: Error Handling — Recovery Patterns", "Retry, fallback, timeout, and graceful degradation")

code_box(s, """\
import asyncio, random
from a2a.client.errors import A2AClientHTTPError
from a2a.client        import A2AClient

async def send_with_retry(client: A2AClient, request,
                          max_retries=3):
    \"\"\"Exponential backoff retry for transient errors.\"\"\"
    retryable = {-32603, -32700}  # Internal, Parse

    for attempt in range(max_retries):
        try:
            resp = await client.send_task(request)
            task = resp.result

            if task.status.state == "failed":
                err_msg = str(task.status.message)
                print(f"Task failed: {err_msg}")
                return None           # do not retry task failures

            return task

        except A2AClientHTTPError as e:
            if e.status_code in (401, 403):
                raise   # auth errors — do not retry

            wait = (2 ** attempt) + random.uniform(0, 1)
            print(f"Attempt {attempt+1} failed "
                  f"(HTTP {e.status_code}), "
                  f"retrying in {wait:.1f}s...")
            await asyncio.sleep(wait)

    return None   # exhausted retries


async def search_with_timeout(client, request,
                               timeout_secs=30.0):
    \"\"\"Cancel task if it exceeds timeout.\"\"\"
    task_id = request.params.id

    try:
        return await asyncio.wait_for(
            client.send_task(request),
            timeout=timeout_secs
        )
    except asyncio.TimeoutError:
        print(f"Timeout — canceling task {task_id}")
        try:
            await client.cancel_task(task_id)
        except Exception:
            pass
        return None
""", 0.3, 1.22, 7.35, 6.0, sz=9.0, label="Retry + Timeout Patterns")

section_box(s, "Resilience Strategies", [
    "Exponential Backoff:",
    "    Retry transient errors (5xx, network)",
    "    Do NOT retry 401/403/task failures",
    "    Add jitter to avoid thundering herd",
    "",
    "Timeout + Cancel:",
    "    Set per-request timeout",
    "    Call tasks/cancel on timeout",
    "    Prevents zombie tasks on server",
    "",
    "Graceful Degradation:",
    "    Flight fails → skip, show error",
    "    Hotel fails → suggest manual search",
    "    Never block entire itinerary",
    "",
    "Circuit Breaker:",
    "    Track consecutive failures per agent",
    "    Stop calling failing agent temporarily",
], 7.8, 1.22, 5.2, 6.0, bsz=11)
notes(s, "Three resilience patterns: (1) exponential backoff for transient server errors — do NOT retry 401/403 or task failures; (2) timeout + cancel to prevent stuck clients; (3) graceful degradation so one agent failure doesn't block the whole orchestration. Never retry auth errors.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — SECTION 11: MULTI-AGENT ORCHESTRATION
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 11: Multi-Agent Orchestration", "Orchestrator as ADK LlmAgent — A2AClient as tool")

code_box(s, """\
from google.adk.agents import LlmAgent
from a2a.client import A2AClient, A2ACardResolver
import asyncio, httpx

class TravelOrchestrator:
    \"\"\"ADK LlmAgent that orchestrates 4 A2A agents.\"\"\"

    def __init__(self):
        self.agents = {}

    async def setup(self):
        \"\"\"Fetch all Agent Cards at startup.\"\"\"
        agent_urls = {
            "flight":  "http://localhost:8001/",
            "hotel":   "http://localhost:8002/",
            "booking": "http://localhost:8003/",
            "weather": "http://localhost:8004/",
        }
        for name, url in agent_urls.items():
            card   = await A2ACardResolver(
                         httpx_client=httpx.AsyncClient(),
                         base_url=url
                     ).get_agent_card()
            client = A2AClient(
                         httpx_client=self._make_client(name),
                         agent_card=card
                     )
            self.agents[name] = client

    async def plan_trip(self, user_query: str):
        \"\"\"Parallel fan-out then sequential booking.\"\"\"
        # Fan out to flight, hotel, weather in parallel
        flight_task, hotel_task, weather_task = (
            await asyncio.gather(
                self._call_agent("flight",  user_query),
                self._call_agent("hotel",   user_query),
                self._call_agent("weather", user_query),
            )
        )

        # Sequential: booking needs flight+hotel data
        booking = await self._book(flight_task, hotel_task)

        return self._aggregate(
            flight_task, hotel_task, weather_task, booking
        )
""", 0.3, 1.22, 7.6, 6.05, sz=9.0, label="TravelOrchestrator — parallel + sequential")

section_box(s, "Orchestration Patterns", [
    "Parallel Fan-Out:",
    "    asyncio.gather() for independent tasks",
    "    Flight + Hotel + Weather run concurrently",
    "    Reduces total latency significantly",
    "",
    "Sequential Chaining:",
    "    Booking depends on flight+hotel results",
    "    Pass previous artifacts as input",
    "",
    "Skill Matching:",
    "    Read skill.description (LLM-readable)",
    "    Match user intent to best agent",
    "    Check inputModes compatibility",
    "",
    "Capability Checking:",
    "    streaming → use sendSubscribe",
    "    pushNotifications → set webhook first",
], 8.1, 1.22, 4.9, 6.05, bsz=11)
notes(s, "The orchestrator is itself an ADK LlmAgent. It fetches Agent Cards at startup, then uses asyncio.gather() to fan out to flight, hotel, and weather agents in parallel. Booking must be sequential because it depends on flight+hotel results. The orchestrator aggregates all results into a unified itinerary.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — ORCHESTRATION: ROUTING AND AGGREGATION
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 11: Orchestration — Routing and Aggregation", "How the orchestrator routes and combines results")

code_box(s, """\
async def _call_agent(self, agent_name: str, query: str):
    \"\"\"Route to named agent with appropriate auth.\"\"\"
    client = self.agents[agent_name]
    card   = client.agent_card

    # Check capability before using streaming
    use_stream = card.capabilities.streaming

    if use_stream:
        return await self._call_streaming(client, query)
    else:
        return await self._call_sync(client, query)

def _aggregate(self, flight, hotel, weather, booking):
    \"\"\"Combine all results into unified itinerary.\"\"\"
    return {
        "itinerary": {
            "flights":  self._extract_data(flight),
            "hotels":   self._extract_data(hotel),
            "weather":  self._extract_data(weather),
            "booking":  self._extract_data(booking),
        },
        "summary": self._generate_summary(
            flight, hotel, weather, booking
        ),
        "totalCost": (
            self._extract_data(flight).get("price", 0) +
            self._extract_data(hotel).get("price", 0)
        )
    }

def _extract_data(self, task):
    \"\"\"Extract DataPart from first artifact.\"\"\"
    if not task or not task.artifacts:
        return {}
    for part in task.artifacts[0].parts:
        if hasattr(part, "data"):
            return part.data
    return {}
""", 0.3, 1.22, 7.35, 6.0, sz=9.0, label="Routing + Aggregation")

# Orchestration flow diagram
rect(s, 7.9, 1.22, 5.1, 6.0, WHITE, LIGHT_GRAY, 1)
rect(s, 7.9, 1.22, 5.1, 0.35, DARK_BLUE)
tb(s, "Execution Flow", 8.0, 1.24, 4.9, 0.3, sz=12, bold=True, color=WHITE)

flow_steps = [
    (GREEN,     "1. Fetch Agent Cards (startup)"),
    (LIGHT_BLUE,"2. Receive user query"),
    (LIGHT_BLUE,"3. Match intent → agents"),
    (ORANGE,    "4. PARALLEL: flight+hotel+weather"),
    (ORANGE,    "5. Await all 3 results"),
    (PURPLE,    "6. SEQUENTIAL: send to booking"),
    (PURPLE,    "7. Handle input-required loop"),
    (GREEN,     "8. Aggregate → unified response"),
    (GREEN,     "9. Return itinerary to user"),
]
for i,(color,text) in enumerate(flow_steps):
    rect(s, 8.05, 1.65+i*0.57, 4.75, 0.47, color, WHITE, 1)
    tb(s, text, 8.15, 1.68+i*0.57, 4.6, 0.38, sz=11, color=WHITE)
notes(s, "The orchestrator checks each agent's streaming capability before calling it. For streaming agents it uses sendSubscribe; for others it uses send. After parallel fan-out it aggregates DataParts from all artifacts into a single unified itinerary response.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — SECTION 12: RUNNING THE SYSTEM — SETUP
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 12: Running the Complete System", "Docker Compose setup and port map")

# Port map table
rect(s, 0.3, 1.22, 12.7, 0.4, DARK_BLUE)
for label, x, w in [("Service",0.4,3.5),("Port",3.95,1.2),("Auth",5.25,2.0),("Features",7.35,5.5)]:
    tb(s, label, x, 1.24, w, 0.32, sz=12, bold=True, color=WHITE)

rows = [
    ("Flight Agent",    "8001","Bearer Token", "Sync tasks, SSE streaming, state history",     GREEN),
    ("Hotel Agent",     "8002","API Key",       "Multimodal PDF input, structured output",      MID_BLUE),
    ("Booking Agent",   "8003","None",          "Input-required, push notifications",           PURPLE),
    ("Weather Agent",   "8004","None",          "SSE streaming, real-time forecasts",           ORANGE),
    ("Travel Orchestrator","8010","None",       "ADK LlmAgent, A2AClient, parallel fan-out",   DARK_BLUE),
    ("Webhook Receiver","9000","Token header",  "FastAPI, validates X-A2A-Notification-Token",  RGBColor(0x5A,0x5A,0x5A)),
]
for i,(svc,port,auth,feat,color) in enumerate(rows):
    bg = WHITE if i%2==0 else RGBColor(0xEE,0xF2,0xF8)
    rect(s, 0.3, 1.64+i*0.52, 12.7, 0.52, bg, LIGHT_GRAY, 1)
    rect(s, 0.3, 1.64+i*0.52, 0.12, 0.52, color)
    tb(s, svc,  0.5,  1.67+i*0.52, 3.35, 0.38, sz=12, bold=True, color=NEAR_BLACK)
    tb(s, port, 3.95, 1.67+i*0.52, 1.1,  0.38, sz=12, color=DARK_BLUE, font="Courier New")
    tb(s, auth, 5.25, 1.67+i*0.52, 1.9,  0.38, sz=11, color=NEAR_BLACK)
    tb(s, feat, 7.35, 1.67+i*0.52, 5.55, 0.38, sz=11, color=NEAR_BLACK)

code_box(s, """\
# Launch everything with Docker Compose
docker-compose up --build

# Watch logs from specific service
docker-compose logs -f flight-agent

# Run the client
python client/main.py "Plan a 5-day trip to Paris in December"
""", 0.3, 5.75, 8.0, 1.65, sz=10.5, label="Quick Start Commands")

section_box(s, "Prerequisites", [
    "Docker + Docker Compose installed",
    "Python 3.11+ with google-a2a-sdk",
    "Set env vars in .env file:",
    "    FLIGHT_BEARER_TOKEN",
    "    HOTEL_API_KEY",
    "    WEBHOOK_TOKEN",
], 8.55, 5.75, 4.45, 1.65, bsz=11)
notes(s, "The complete system runs via docker-compose up. Six services total: 4 specialist agents, 1 orchestrator, 1 webhook receiver. Each agent is a separate FastAPI app. The client sends a natural language query and receives a fully planned itinerary.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — DOCKER COMPOSE AND AGENT SERVER SETUP
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 12: Docker Compose Configuration", "Service definitions for all 6 components")

code_box(s, """\
# docker-compose.yml
version: "3.9"
services:
  flight-agent:
    build: ./agents/flight
    ports: ["8001:8001"]
    environment:
      - BEARER_TOKEN=${FLIGHT_BEARER_TOKEN}
    command: uvicorn app:app --host 0.0.0.0 --port 8001

  hotel-agent:
    build: ./agents/hotel
    ports: ["8002:8002"]
    environment:
      - API_KEY=${HOTEL_API_KEY}
    command: uvicorn app:app --host 0.0.0.0 --port 8002

  booking-agent:
    build: ./agents/booking
    ports: ["8003:8003"]
    environment:
      - WEBHOOK_TOKEN=${WEBHOOK_TOKEN}
    command: uvicorn app:app --host 0.0.0.0 --port 8003

  weather-agent:
    build: ./agents/weather
    ports: ["8004:8004"]
    command: uvicorn app:app --host 0.0.0.0 --port 8004

  orchestrator:
    build: ./orchestrator
    ports: ["8010:8010"]
    depends_on:
      - flight-agent
      - hotel-agent
      - booking-agent
      - weather-agent
    command: uvicorn app:app --host 0.0.0.0 --port 8010

  webhook-receiver:
    build: ./webhook
    ports: ["9000:9000"]
    environment:
      - WEBHOOK_TOKEN=${WEBHOOK_TOKEN}
    command: uvicorn app:app --host 0.0.0.0 --port 9000
""", 0.3, 1.22, 7.35, 6.0, sz=9.0, label="docker-compose.yml")

code_box(s, """\
# agents/flight/app.py — A2A FastAPI server
from a2a.server.apps import A2AStarletteApplication
from a2a.server.request_handlers import DefaultRequestHandler
from a2a.server.tasks import InMemoryTaskStore
from a2a.types import AgentCard, AgentCapabilities

from flight_executor import FlightAgentExecutor

agent_card = AgentCard(
    name="Flight Search Agent",
    url="http://localhost:8001/",
    version="1.0.0",
    capabilities=AgentCapabilities(
        streaming=True,
        pushNotifications=False
    ),
    ...
)
handler = DefaultRequestHandler(
    agent_executor=FlightAgentExecutor(),
    task_store=InMemoryTaskStore(),
)
app = A2AStarletteApplication(
    agent_card=agent_card,
    http_handler=handler
).build()
""", 0.3, 7.28, 7.35, 2.5, sz=9.0, label="Agent Server Boilerplate")

section_box(s, "Project Structure", [
    "a2a-travel/",
    "    agents/",
    "        flight/   (app.py, executor.py)",
    "        hotel/    (app.py, executor.py)",
    "        booking/  (app.py, executor.py)",
    "        weather/  (app.py, executor.py)",
    "    orchestrator/ (app.py, orchestrator.py)",
    "    webhook/      (app.py)",
    "    client/       (main.py)",
    "    docker-compose.yml",
    "    .env.example",
    "    requirements.txt",
], 7.9, 1.22, 5.1, 6.0, bsz=11)
notes(s, "Each agent is a standalone FastAPI app built with A2AStarletteApplication. The framework automatically serves the Agent Card at /.well-known/agent.json and routes JSON-RPC requests to your AgentExecutor. Docker Compose orchestrates startup order with depends_on.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — END-TO-END FLOW
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 12: End-to-End Request Flow", "Tracing a user query through the entire system")

# Flow steps
steps = [
    ("1","User sends query","python client/main.py \"Plan 5-day Paris trip\"","Client",DARK_BLUE),
    ("2","Orchestrator receives","ADK LlmAgent parses intent, identifies 3 subtasks","Orchestrator",LIGHT_BLUE),
    ("3","Fetch Agent Cards","GET /.well-known/agent.json × 4 agents","Discovery",GREEN),
    ("4","Parallel fan-out","asyncio.gather(flight, hotel, weather) → 3 concurrent tasks","A2A Tasks",ORANGE),
    ("5","Flight streams","tasks/sendSubscribe → SSE events → collect artifact","Streaming",GREEN),
    ("6","Hotel multimodal","tasks/send + PDF brochure → structured hotel list","Multimodal",MID_BLUE),
    ("7","Weather streams","tasks/sendSubscribe → SSE → 5-day forecast artifact","Streaming",ORANGE),
    ("8","Booking: input loop","tasks/send → input-required × 3 → completed","Input Loop",PURPLE),
    ("9","Aggregate","Combine all 4 artifacts → unified itinerary dict","Aggregate",GREEN),
    ("10","Return to user","Print formatted travel plan with booking ref","Response",DARK_BLUE),
]
rect(s, 0.2, 1.22, 12.9, 0.35, DARK_BLUE)
for label, x, w in [("Step",0.3,0.6),("Action",0.95,2.5),("Details",3.55,6.5),("Phase",10.15,1.2),]:
    tb(s, label, x, 1.24, w, 0.28, sz=10, bold=True, color=WHITE)

for i,(num,action,detail,phase,color) in enumerate(steps):
    bg = WHITE if i%2==0 else RGBColor(0xEE,0xF2,0xF8)
    rect(s, 0.2, 1.59+i*0.52, 12.9, 0.52, bg, LIGHT_GRAY, 1)
    rect(s, 0.2, 1.59+i*0.52, 0.12, 0.52, color)
    tb(s, num,    0.38, 1.62+i*0.52, 0.55, 0.38, sz=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    tb(s, action, 0.98, 1.62+i*0.52, 2.45, 0.38, sz=11, bold=True, color=NEAR_BLACK)
    tb(s, detail, 3.52, 1.62+i*0.52, 6.5,  0.38, sz=10, color=NEAR_BLACK)
    rect(s, 10.05, 1.64+i*0.52, 3.0, 0.38, color)
    tb(s, phase, 10.05, 1.64+i*0.52, 3.0, 0.38, sz=10, color=WHITE, align=PP_ALIGN.CENTER)
notes(s, "End-to-end trace: query enters at client, orchestrator routes to 4 agents in parallel/sequential mix, results aggregated into final itinerary. Steps 4-7 run concurrently via asyncio.gather(). Step 8 is sequential and may loop 3 times for input-required.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — FEATURE COVERAGE TABLE
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Complete A2A Feature Coverage", "Every protocol feature demonstrated in this tutorial")

features = [
    ("Agent Cards (/.well-known/agent.json)","Section 2","All 4 agents","✓"),
    ("tasks/send (synchronous)","Section 3","Hotel, Booking","✓"),
    ("tasks/sendSubscribe (SSE streaming)","Section 5","Flight, Weather","✓"),
    ("tasks/get (polling)","Section 4","All agents","✓"),
    ("tasks/cancel","Section 4","All agents","✓"),
    ("TextPart","Section 3","All agents","✓"),
    ("FilePart — inline bytes (PDF)","Section 6","Hotel Agent","✓"),
    ("FilePart — hosted URI","Section 6","Hotel Agent","✓"),
    ("DataPart — structured JSON","Section 3,6","Hotel, Booking","✓"),
    ("Task state: input-required","Section 7","Booking Agent","✓"),
    ("Push notifications (webhook)","Section 8","Booking Agent","✓"),
    ("tasks/pushNotification/set","Section 8","Booking Agent","✓"),
    ("Bearer Token authentication","Section 9","Flight Agent","✓"),
    ("API Key authentication","Section 9","Hotel Agent","✓"),
    ("capabilities.streaming: true","Section 5","Flight, Weather","✓"),
    ("capabilities.pushNotifications: true","Section 8","Booking Agent","✓"),
    ("stateTransitionHistory","Section 4","Flight Agent","✓"),
    ("Parallel fan-out (asyncio.gather)","Section 11","Orchestrator","✓"),
    ("Sequential chaining","Section 11","Orchestrator","✓"),
    ("Error codes + retry patterns","Section 10","Client","✓"),
]

rect(s, 0.2, 1.22, 12.9, 0.35, DARK_BLUE)
for label, x, w in [("A2A Feature",0.3,5.5),("Section(s)",5.9,2.0),("Agent/Component",8.0,3.0),("Status",11.1,1.9)]:
    tb(s, label, x, 1.24, w, 0.28, sz=10, bold=True, color=WHITE)

for i,(feat,sect,agent,status) in enumerate(features):
    bg = WHITE if i%2==0 else RGBColor(0xEE,0xF2,0xF8)
    rect(s, 0.2, 1.59+i*0.285, 12.9, 0.285, bg, LIGHT_GRAY, 1)
    tb(s, feat,   0.3,  1.6+i*0.285, 5.45, 0.25, sz=9.5, color=NEAR_BLACK)
    tb(s, sect,   5.9,  1.6+i*0.285, 1.9,  0.25, sz=9.5, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    tb(s, agent,  8.0,  1.6+i*0.285, 2.9,  0.25, sz=9.5, color=NEAR_BLACK)
    tb(s, status, 11.1, 1.6+i*0.285, 1.85, 0.25, sz=9.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
notes(s, "Complete feature coverage matrix. Every A2A protocol feature from the specification is demonstrated in this tutorial: all Part types, all task states, both streaming and sync patterns, both auth schemes, push notifications, and multi-agent orchestration patterns.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — SAMPLE CLIENT OUTPUT
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Section 12: Sample Client Run", "What you see when running client/main.py")

code_box(s, """\
$ python client/main.py "Plan a 5-day trip to Paris in December"

[Orchestrator] Fetching agent cards...
  ✓ Flight Agent  (streaming=True,  auth=Bearer)
  ✓ Hotel Agent   (streaming=False, auth=ApiKey)
  ✓ Booking Agent (streaming=False, auth=None, push=True)
  ✓ Weather Agent (streaming=True,  auth=None)

[Orchestrator] Fanning out to 3 agents in parallel...

[Flight  → streaming] AF 001 NYC→CDG  Dec 15  $842  (14h 05m) ▌
[Flight  → streaming] UA 088 NYC→CDG  Dec 15  $919  (11h 45m) ▌
[Weather → streaming] Paris Dec 15-20: 8°C, light rain, wind 12km/h ▌

[Hotel] Sent PDF brochure. Received structured hotel list.
  - Le Marais Boutique Hotel: €145/night, 4.6★, central Paris
  - Opera Grand Hotel:        €210/night, 4.8★, near Louvre

[Orchestrator] Sending to Booking Agent...
[Booking] Status: input-required
  Agent: "Which seat class? (economy / business / first)"
  → economy

[Booking] Status: input-required
  Agent: "Meal preference? (standard / vegetarian / vegan)"
  → vegetarian

[Booking] Status: input-required
  Agent: "Loyalty number? (enter number or 'none')"
  → none

[Booking] Status: completed ✓

━━━━━━━━━━━━━━━ TRAVEL ITINERARY ━━━━━━━━━━━━━━━
Flight:  AF 001 NYC→CDG Dec 15 | Economy | $842
Hotel:   Le Marais Boutique Hotel | 5 nights | €725
Weather: 8°C, light rain — pack layers and umbrella
Booking: CONF-2024-PAR-00142 | Total: ~$1,600
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
""", 0.3, 1.22, 12.7, 6.05, sz=9.5, label="Terminal Output — python client/main.py")
notes(s, "Sample run shows the full flow: agent card discovery, parallel fan-out with streaming output, hotel multimodal response, the 3-step booking input-required loop, and the final aggregated itinerary with booking confirmation reference.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — AGENT EXECUTOR PATTERN (SERVER)
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Deep Dive: AgentExecutor Pattern", "How to implement any A2A-compliant agent")

code_box(s, """\
from a2a.server.agent_execution import (
    AgentExecutor, RequestContext
)
from a2a.server.events import EventQueue
from a2a.types import (
    TaskStatus, TaskStatusUpdateEvent,
    TaskArtifactUpdateEvent, Artifact, DataPart
)

class HotelAgentExecutor(AgentExecutor):
    \"\"\"Hotel Agent: accepts text + PDF, returns JSON.\"\"\"

    async def execute(self,
                      context: RequestContext,
                      event_queue: EventQueue) -> None:

        # ── Extract parts from request ─────────────────
        text_query = ""
        pdf_bytes  = None

        for part in context.message.parts:
            if hasattr(part, "text"):
                text_query = part.text
            elif hasattr(part, "file"):
                if part.file.mimeType == "application/pdf":
                    pdf_bytes = part.file.bytes

        # ── Signal working ─────────────────────────────
        await event_queue.enqueue_event(
            TaskStatusUpdateEvent(
                status=TaskStatus(state="working"),
                final=False
            )
        )

        # ── Do actual work ─────────────────────────────
        hotels = await self._search_hotels(
            query=text_query,
            brochure_pdf=pdf_bytes
        )

        # ── Return structured result ───────────────────
        await event_queue.enqueue_event(
            TaskArtifactUpdateEvent(
                artifact=Artifact(
                    name="hotel_results",
                    index=0,
                    lastChunk=True,
                    parts=[DataPart(type="data", data=hotels)]
                )
            )
        )

        # ── Signal completion ──────────────────────────
        await event_queue.enqueue_event(
            TaskStatusUpdateEvent(
                status=TaskStatus(state="completed"),
                final=True
            )
        )
""", 0.3, 1.22, 7.6, 6.05, sz=8.8, label="HotelAgentExecutor — complete implementation")

section_box(s, "AgentExecutor Interface", [
    "async execute(context, event_queue)",
    "    context.message — incoming Message",
    "    context.task_id — current task ID",
    "    context.session_id — session ID",
    "    event_queue.enqueue_event() — emit",
    "",
    "Always emit in this order:",
    "    1. working status (non-final)",
    "    2. artifact event(s) with content",
    "    3. completed/failed status (final=True)",
    "",
    "For input-required:",
    "    Emit input-required status",
    "    Include question in status.message",
    "    Client resumes with same task_id",
    "",
    "Framework handles:",
    "    JSON-RPC wrapping",
    "    SSE formatting",
    "    Task store (InMemoryTaskStore)",
    "    Agent Card serving",
], 8.1, 1.22, 4.9, 6.05, bsz=10)
notes(s, "The AgentExecutor is the core pattern for building A2A agents. Implement the async execute() method. Always emit a working status first, then one or more artifact events, then a final completed or failed status. The framework handles all protocol details automatically.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — BOOKING AGENT: INPUT REQUIRED SERVER
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Deep Dive: Booking Agent — Input-Required Server", "Multi-step conversation state on the server side")

code_box(s, """\
from a2a.server.agent_execution import AgentExecutor
from a2a.types import (
    TaskStatus, TaskStatusUpdateEvent, TaskArtifactUpdateEvent,
    Artifact, DataPart, Message, TextPart
)

class BookingAgentExecutor(AgentExecutor):

    # In-memory session state (use Redis in production)
    _sessions: dict = {}

    async def execute(self, context, event_queue):
        session_id = context.session_id
        task_id    = context.task_id
        state      = self._sessions.get(session_id, {})

        user_text = self._get_text(context.message)

        # ── Step 1: Need seat class ────────────────────
        if "seat_class" not in state:
            state["booking_data"] = self._get_data(context.message)
            self._sessions[session_id] = state
            await event_queue.enqueue_event(
                TaskStatusUpdateEvent(
                    status=TaskStatus(
                        state="input-required",
                        message=Message(role="agent", parts=[
                            TextPart(type="text",
                                     text="Which seat class? "
                                          "(economy/business/first)")
                        ])
                    ), final=True
                )
            )
            return

        # ── Step 2: Need meal preference ──────────────
        if "seat_class" in state and "meal" not in state:
            state["seat_class"] = user_text
            self._sessions[session_id] = state
            await event_queue.enqueue_event(
                TaskStatusUpdateEvent(
                    status=TaskStatus(
                        state="input-required",
                        message=Message(role="agent", parts=[
                            TextPart(type="text",
                                     text="Meal preference? "
                                          "(standard/vegetarian/vegan)")
                        ])
                    ), final=True
                )
            )
            return

        # ── (additional steps omitted for brevity) ────
        # Step 3: loyalty number, then complete booking
        booking_ref = await self._confirm_booking(state)
        await event_queue.enqueue_event(
            TaskArtifactUpdateEvent(
                artifact=Artifact(name="booking", index=0,
                    lastChunk=True,
                    parts=[DataPart(type="data",
                                    data=booking_ref)])
            )
        )
        await event_queue.enqueue_event(
            TaskStatusUpdateEvent(
                status=TaskStatus(state="completed"),
                final=True
            )
        )
        del self._sessions[session_id]
""", 0.3, 1.22, 7.8, 6.05, sz=8.5, label="BookingAgentExecutor — input-required server")

section_box(s, "Server-Side State Machine", [
    "Use session_id as state key",
    "Store partial answers in _sessions dict",
    "Check what data is missing each call",
    "Emit input-required with question message",
    "Set final=True to end current SSE stream",
    "Client reuses same task_id to resume",
    "",
    "Production Considerations:",
    "    Use Redis for _sessions (not in-memory)",
    "    Set TTL on session state (e.g. 30 min)",
    "    Handle task_id collisions gracefully",
    "    Validate each answer before proceeding",
    "    Sanitize user input before using",
], 8.15, 1.22, 4.9, 6.05, bsz=11)
notes(s, "Server-side input-required is a state machine keyed by session_id. On each call check what information you already have and ask for what is missing. Set final=True on input-required status events since each call ends the current stream. The client resumes by sending to the same task_id.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 29 — SESSION ID AND HISTORY
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Deep Dive: SessionId and Task History", "Grouping tasks and accessing conversation history")

code_box(s, """\
import uuid
from a2a.types import (
    GetTaskRequest, TaskQueryParams,
    SendTaskRequest, TaskSendParams, Message, TextPart
)

# ── Session groups related tasks ──────────────────────
SESSION_ID = f"travel-session-{uuid.uuid4()}"

# ── Task 1: Search flights ────────────────────────────
flight_task_id = str(uuid.uuid4())
flight_resp = await flight_client.send_task(
    SendTaskRequest(id=flight_task_id, params=TaskSendParams(
        id=flight_task_id,
        sessionId=SESSION_ID,
        message=Message(role="user", parts=[
            TextPart(type="text", text="NYC to Paris Dec 15")
        ])
    ))
)

# ── Task 2: Search hotels (same session) ─────────────
hotel_task_id = str(uuid.uuid4())
hotel_resp = await hotel_client.send_task(
    SendTaskRequest(id=hotel_task_id, params=TaskSendParams(
        id=hotel_task_id,
        sessionId=SESSION_ID,           # same session!
        message=Message(role="user", parts=[
            TextPart(type="text", text="Hotels in Paris Dec 15-20")
        ])
    ))
)

# ── Retrieve task with history ─────────────────────────
detailed = await flight_client.get_task(
    GetTaskRequest(id=flight_task_id, params=TaskQueryParams(
        id=flight_task_id,
        historyLength=10    # return last 10 message turns
    ))
)

# task.history is list[Message]
for msg in detailed.result.history:
    print(f"[{msg.role}] {msg.parts[0].text[:80]}")
""", 0.3, 1.22, 7.5, 6.05, sz=9.0, label="SessionId Grouping + Task History")

section_box(s, "SessionId Semantics", [
    "Client-generated string (any format)",
    "Groups logically related tasks",
    "Passed to every task in a conversation",
    "Agents can share session context",
    "Useful for multi-turn dialogues",
    "",
    "Task History (historyLength):",
    "    Returned via tasks/get",
    "    List of Message objects",
    "    Most recent N turns",
    "    Requires stateTransitionHistory:true",
    "    in Agent Card capabilities",
    "",
    "History vs Artifacts:",
    "    history — conversation messages",
    "    artifacts — deliverable outputs",
], 8.0, 1.22, 5.0, 6.05, bsz=11)
notes(s, "SessionId groups related tasks together. It's client-generated and passed on every request. When you call tasks/get with historyLength > 0 you get the conversation history. This requires the agent to advertise stateTransitionHistory:true in its capabilities. Flight Agent enables this for debugging.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 30 — SUMMARY AND KEY TAKEAWAYS
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, BODY_BG)
title_bar(s, "Summary and Key Takeaways", "What you've learned in this tutorial")

takeaways = [
    (DARK_BLUE, "A2A Protocol",
     ["Open HTTP + JSON-RPC 2.0 standard",
      "Framework-agnostic: works with any agent framework",
      "Agent Cards enable automatic discovery"]),
    (GREEN, "Tasks & Lifecycle",
     ["Every interaction is a Task with a unique UUID",
      "6 states: submitted→working→completed/failed/canceled/input-required",
      "Poll with tasks/get or stream with tasks/sendSubscribe"]),
    (MID_BLUE, "Content Types",
     ["TextPart: plain text",
      "FilePart: inline bytes (base64) or hosted URI",
      "DataPart: structured JSON for machine consumption"]),
    (ORANGE, "Advanced Features",
     ["SSE streaming: real-time chunked output",
      "Input-required: human-in-the-loop multi-step",
      "Push notifications: webhook for long tasks"]),
    (PURPLE, "Security",
     ["Bearer tokens and API Keys declared in Agent Card",
      "Clients auto-discover and inject credentials",
      "Always validate on every server request"]),
    (RGBColor(0x5A,0x5A,0x5A), "Orchestration",
     ["Parallel fan-out with asyncio.gather()",
      "Sequential chaining for dependent tasks",
      "Graceful degradation on agent failure"]),
]
for i,(color,title,bullets) in enumerate(takeaways):
    col = i % 2
    row = i // 2
    x = 0.3 if col == 0 else 6.8
    y = 1.25 + row * 2.0
    rect(s, x, y, 6.2, 1.85, WHITE, LIGHT_GRAY, 1)
    rect(s, x, y, 6.2, 0.38, color)
    tb(s, title, x+0.12, y+0.05, 5.9, 0.3, sz=13, bold=True, color=WHITE)
    multi_para(s, bullets, x+0.15, y+0.44, 5.9, 1.35, sz=12, color=NEAR_BLACK)
notes(s, "Summary slide covering all 6 major areas: A2A protocol fundamentals, task lifecycle, content types, advanced features (streaming/input-required/push), security, and orchestration patterns. You now have everything needed to build production A2A multi-agent systems.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 31 — NEXT STEPS AND RESOURCES
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide()
set_bg(s, DARK_BLUE)
rect(s, 0, 0, 0.5, 7.5, ACCENT)
title_bar(s, "Next Steps and Resources", "Where to go from here")

section_box(s, "Official Resources", [
    "A2A Spec:    github.com/google/A2A",
    "ADK Docs:    google.github.io/adk-docs",
    "ADK Samples: github.com/google/adk-samples",
    "A2A Python:  pip install google-a2a",
    "PyPI:        pypi.org/project/google-a2a",
], 0.7, 1.22, 5.6, 2.5, bsz=12)

section_box(s, "This Tutorial's Code", [
    "All 4 agent executors (Flight/Hotel/Booking/Weather)",
    "TravelOrchestrator with parallel fan-out",
    "Webhook receiver (FastAPI)",
    "Docker Compose for full system",
    "Client with input-required loop handling",
    "Complete error handling and retry logic",
], 6.55, 1.22, 6.45, 2.5, bsz=12)

section_box(s, "Extend This Project", [
    "Add a Car Rental agent (new A2A skill)",
    "Add OAuth2 authentication to Booking Agent",
    "Replace InMemoryTaskStore with Redis",
    "Add OpenTelemetry tracing across agents",
    "Build a web UI on top of the orchestrator",
    "Deploy to Cloud Run with real APIs",
    "Add a recommendation agent using RAG",
], 0.7, 3.85, 5.6, 3.0, bsz=12)

section_box(s, "Key A2A Best Practices", [
    "Always validate Agent Cards before sending",
    "Use asyncio.gather() for parallel tasks",
    "Handle ALL task states including input-required",
    "Set timeouts and cancel on timeout",
    "Use DataPart for machine-readable outputs",
    "Declare capabilities accurately in Agent Card",
    "Never store sensitive data in task history",
    "Use sessionId for multi-turn conversations",
], 6.55, 3.85, 6.45, 3.0, bsz=12)
notes(s, "Resources: the official A2A specification is on GitHub at google/A2A. Google ADK documentation and samples are also on GitHub. Install the Python SDK with: pip install google-a2a. The tutorial code is self-contained and runnable with Docker Compose.")

# ════════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════════
output_path = "/home/agenticai/dev/a2a/A2A_Tutorial.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
