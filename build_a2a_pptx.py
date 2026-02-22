"""
Build A2A_Protocol_Concepts.pptx — 20 slides, 16:9 widescreen
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x00, 0x30, 0x87)
GOLD   = RGBColor(0xFF, 0xB3, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF4, 0xF6, 0xF9)
DGRAY  = RGBColor(0x1E, 0x1E, 0x1E)
BLACK  = RGBColor(0x00, 0x00, 0x00)
MGRAY  = RGBColor(0x60, 0x60, 0x60)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
TEAL   = RGBColor(0x00, 0x69, 0x6E)
ORANGE = RGBColor(0xE6, 0x5C, 0x00)
BLUE2  = RGBColor(0x15, 0x65, 0xC0)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, l, t, w, h, text, size=14, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, wrap=True, italic=False, font_name="Calibri"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font_name
    return txb


def header_bar(slide, title, subtitle=None):
    """Navy header bar across the top, gold left accent stripe."""
    # Gold left stripe
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, fill=GOLD)
    # Navy header
    add_rect(slide, Inches(0.12), Inches(0), W - Inches(0.12), Inches(1.1), fill=NAVY)
    # Title text
    add_textbox(slide, Inches(0.35), Inches(0.1), Inches(12.5), Inches(0.65),
                title, size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, Inches(0.35), Inches(0.72), Inches(12.5), Inches(0.32),
                    subtitle, size=13, bold=False, color=GOLD, align=PP_ALIGN.LEFT)
    # Light gray body background
    add_rect(slide, Inches(0.12), Inches(1.1), W - Inches(0.12), H - Inches(1.1), fill=LGRAY)


def bullet_box(slide, l, t, w, h, bullets, title=None, title_color=NAVY,
               bullet_size=13, title_size=15, bg=WHITE, border=NAVY):
    add_rect(slide, l, t, w, h, fill=bg, line=border, line_w=Pt(1.2))
    y = t + Inches(0.12)
    if title:
        add_textbox(slide, l + Inches(0.12), y, w - Inches(0.24), Inches(0.36),
                    title, size=title_size, bold=True, color=title_color)
        y += Inches(0.38)
    for b in bullets:
        add_textbox(slide, l + Inches(0.18), y, w - Inches(0.32), Inches(0.3),
                    b, size=bullet_size, color=BLACK)
        y += Inches(0.29)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width=Pt(2)):
    """Add a simple line connector (arrow-like)."""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # STRAIGHT=1
    connector.line.color.rgb = color
    connector.line.width = width


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# Full navy background
add_rect(sl, 0, 0, W, H, fill=NAVY)
# Gold accent left stripe
add_rect(sl, 0, 0, Inches(0.5), H, fill=GOLD)
# Gold bottom stripe
add_rect(sl, 0, H - Inches(0.35), W, Inches(0.35), fill=GOLD)

# Main title
add_textbox(sl, Inches(1), Inches(1.6), Inches(11), Inches(1.2),
            "Agent2Agent (A2A) Protocol",
            size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_textbox(sl, Inches(1), Inches(2.85), Inches(11), Inches(0.6),
            "An Open Standard for AI Agent Interoperability",
            size=24, bold=False, color=GOLD, align=PP_ALIGN.LEFT)

# Divider line
add_rect(sl, Inches(1), Inches(3.55), Inches(10), Inches(0.04), fill=GOLD)

# Origin / governance
add_textbox(sl, Inches(1), Inches(3.75), Inches(11), Inches(0.4),
            "Originally created by Google  |  Now managed by the Linux Foundation",
            size=15, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

# Date
add_textbox(sl, Inches(1), Inches(4.25), Inches(4), Inches(0.35),
            "2026", size=14, bold=False, color=GOLD)

# Tagline
add_textbox(sl, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
            "\"A common language for AI agents across frameworks, vendors, and clouds\"",
            size=16, bold=False, italic=True, color=RGBColor(0xB0, 0xC4, 0xDE),
            align=PP_ALIGN.LEFT)

add_notes(sl, "Welcome slide. A2A is an open protocol enabling AI agents built on different frameworks "
              "to discover, communicate, and collaborate. Originally developed by Google, now stewarded "
              "by the Linux Foundation for open governance.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Agenda", "What we'll cover today")

topics_left = [
    "1.  The Problem — Before A2A",
    "2.  What is A2A?",
    "3.  A2A in the AI Stack",
    "4.  Key Design Principles",
    "5.  Core Actors",
    "6.  Fundamental Concepts",
    "7.  Agent Card Deep Dive",
    "8.  The Task Lifecycle",
    "9.  Messages & Parts",
    "10. Interaction Patterns",
]
topics_right = [
    "11. The 11 Core Operations",
    "12. Protocol Bindings",
    "13. Authentication & Security",
    "14. Multi-Turn Interactions",
    "15. Request Lifecycle",
    "16. SDK & Framework Support",
    "17. Key Takeaways",
    "18. Resources",
    "",
    "",
]

# Two-column layout
col_w = Inches(5.9)
col_h = Inches(5.6)
for col_idx, topics in enumerate([topics_left, topics_right]):
    lx = Inches(0.35 + col_idx * 6.45)
    add_rect(sl, lx, Inches(1.2), col_w, col_h, fill=WHITE, line=NAVY, line_w=Pt(1))
    y = Inches(1.4)
    for t in topics:
        if t:
            add_textbox(sl, lx + Inches(0.2), y, col_w - Inches(0.3), Inches(0.42),
                        t, size=13.5, color=BLACK)
        y += Inches(0.48)

add_notes(sl, "Overview of all 18 topic slides in this deck. The presentation covers the full A2A "
              "protocol from motivation through specification details and SDK support.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Problem
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "The Problem", "Before A2A: the integration mess")

# BEFORE box
bx = Inches(0.35)
by = Inches(1.25)
bw = Inches(5.8)
bh = Inches(5.8)
add_rect(sl, bx, by, bw, bh, fill=RGBColor(0xFF, 0xEB, 0xEB), line=RGBColor(0xC6, 0x28, 0x28), line_w=Pt(2))
add_textbox(sl, bx + Inches(0.15), by + Inches(0.1), bw - Inches(0.3), Inches(0.4),
            "BEFORE A2A", size=16, bold=True, color=RGBColor(0xC6, 0x28, 0x28))

before_items = [
    "Agents wrapped as 'tools' — loses",
    "  agent negotiation capabilities",
    "",
    "Custom point-to-point integrations",
    "  for every agent pair",
    "",
    "Engineering overhead grows O(n²)",
    "  with agent network size",
    "",
    "No standard security model —",
    "  ad-hoc auth & gaps",
    "",
    "Vendor & framework lock-in",
    "",
    "Example: Trip planning requires",
    "  flight + hotel + currency + tour",
    "  agents — no shared protocol!",
]
y = by + Inches(0.6)
for item in before_items:
    if item:
        add_textbox(sl, bx + Inches(0.25), y, bw - Inches(0.4), Inches(0.3),
                    item, size=12, color=BLACK)
    y += Inches(0.29)

# AFTER box
ax = Inches(6.85)
ay = Inches(1.25)
aw = Inches(6.1)
ah = Inches(5.8)
add_rect(sl, ax, ay, aw, ah, fill=RGBColor(0xE8, 0xF5, 0xE9), line=GREEN, line_w=Pt(2))
add_textbox(sl, ax + Inches(0.15), ay + Inches(0.1), aw - Inches(0.3), Inches(0.4),
            "WITH A2A", size=16, bold=True, color=GREEN)

after_items = [
    "Agents communicate natively via",
    "  open standard protocol",
    "",
    "Single integration surface —",
    "  any A2A agent works with any other",
    "",
    "Scales linearly as agent networks grow",
    "",
    "Built-in enterprise security:",
    "  OAuth 2.0, mTLS, OpenID Connect",
    "",
    "Vendor & framework agnostic",
    "",
    "Trip planning: each specialist agent",
    "  discovered, authenticated, and",
    "  coordinated automatically via A2A",
]
y = ay + Inches(0.6)
for item in after_items:
    if item:
        add_textbox(sl, ax + Inches(0.25), y, aw - Inches(0.4), Inches(0.3),
                    item, size=12, color=BLACK)
    y += Inches(0.29)

# Arrow between boxes
add_rect(sl, Inches(6.15), Inches(3.8), Inches(0.65), Inches(0.4), fill=GOLD)
add_textbox(sl, Inches(6.15), Inches(3.83), Inches(0.65), Inches(0.35),
            "→", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_notes(sl, "The core problem: without a standard protocol, AI agent integration requires custom "
              "point-to-point connections that scale poorly. Wrapping agents as tools loses the "
              "native negotiation and multi-turn reasoning capabilities agents provide. A2A solves "
              "all of these issues through standardization.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What is A2A?
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "What is A2A?", "Definition, origin, and purpose")

# Large tagline
add_rect(sl, Inches(0.35), Inches(1.2), Inches(12.6), Inches(0.85), fill=NAVY)
add_textbox(sl, Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.65),
            "\"A common language for AI agents — enabling seamless collaboration across frameworks, vendors, and clouds\"",
            size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Info cards row
cards = [
    ("What",   NAVY, WHITE,
     ["Open protocol standard for", "AI agent-to-agent communication", "", "Defines how agents discover,", "authenticate, and exchange", "tasks and messages"]),
    ("Origin", TEAL, WHITE,
     ["Created by Google", "Now governed by the", "Linux Foundation", "", "Open governance ensures", "vendor-neutral evolution"]),
    ("Purpose", BLUE2, WHITE,
     ["Enable interoperability across", "diverse agent frameworks", "", "Break down silos between", "vendor ecosystems", "and platforms"]),
    ("Scope",  GREEN, WHITE,
     ["Agent-to-agent (not tool-to-agent)", "", "Complements MCP:", "MCP = agent → tool", "A2A = agent → agent", ""]),
]

cx = Inches(0.35)
cy = Inches(2.2)
cw = Inches(3.05)
ch = Inches(4.6)
gap = Inches(0.12)

for label, bg, fg, lines in cards:
    add_rect(sl, cx, cy, cw, ch, fill=bg, line=None)
    add_textbox(sl, cx + Inches(0.12), cy + Inches(0.12), cw - Inches(0.24), Inches(0.45),
                label, size=18, bold=True, color=GOLD if bg == NAVY else WHITE)
    y2 = cy + Inches(0.65)
    for line in lines:
        add_textbox(sl, cx + Inches(0.15), y2, cw - Inches(0.3), Inches(0.35),
                    line, size=12.5, color=WHITE if fg == WHITE else BLACK)
        y2 += Inches(0.38)
    cx += cw + gap

add_notes(sl, "A2A is an open protocol standard that gives AI agents a common language for "
              "communication. Born at Google and now managed by the Linux Foundation, it enables "
              "any agent — regardless of the underlying framework or vendor — to discover and "
              "collaborate with any other A2A-compliant agent.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — A2A in the AI Stack
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "A2A in the AI Stack", "Where A2A fits alongside MCP and ADK")

# Three layers
layers = [
    (Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.3),
     RGBColor(0x4A, 0x14, 0x8C), WHITE,
     "LAYER 3 — Agent Framework", "ADK · LangChain · CrewAI · Semantic Kernel · AutoGen · Custom"),
    (Inches(1.0), Inches(4.05), Inches(11.3), Inches(1.3),
     NAVY, GOLD,
     "LAYER 2 — Agent ↔ Agent Communication  →  A2A Protocol",
     "Enables multi-turn reasoning · negotiation · task delegation across agents"),
    (Inches(1.0), Inches(2.6), Inches(11.3), Inches(1.3),
     TEAL, WHITE,
     "LAYER 1 — Tools & Data  →  MCP (Model Context Protocol)",
     "Connects models to databases, APIs, file systems, external tools"),
]

for lx, ly, lw, lh, bg, fg, title, sub in layers:
    add_rect(sl, lx, ly, lw, lh, fill=bg)
    add_textbox(sl, lx + Inches(0.2), ly + Inches(0.12), lw - Inches(0.4), Inches(0.45),
                title, size=15, bold=True, color=fg)
    add_textbox(sl, lx + Inches(0.2), ly + Inches(0.6), lw - Inches(0.4), Inches(0.55),
                sub, size=12.5, color=fg)

# Arrows between layers
for arrow_y in [Inches(3.98), Inches(5.45)]:
    add_rect(sl, Inches(6.4), arrow_y - Inches(0.08), Inches(0.5), Inches(0.22), fill=GOLD)
    add_textbox(sl, Inches(6.35), arrow_y - Inches(0.1), Inches(0.6), Inches(0.25),
                "▲", size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Comparison note
add_rect(sl, Inches(1.0), Inches(1.15), Inches(11.3), Inches(1.3), fill=WHITE, line=NAVY, line_w=Pt(1))
cols = [
    (Inches(1.15), "MCP — Model Context Protocol",
     "Connects AI models to tools, databases, and data sources"),
    (Inches(5.2), "A2A — Agent2Agent Protocol",
     "Enables agents to collaborate, delegate, and coordinate with each other"),
    (Inches(9.2), "ADK — Agent Development Kit",
     "Framework for building agents; A2A is the protocol enabling interop"),
]
for cx2, title2, desc2 in cols:
    add_textbox(sl, cx2, Inches(1.2), Inches(3.8), Inches(0.35),
                title2, size=12, bold=True, color=NAVY)
    add_textbox(sl, cx2, Inches(1.55), Inches(3.8), Inches(0.5),
                desc2, size=11, color=MGRAY)

add_notes(sl, "The AI stack has three layers. At the bottom, MCP (Model Context Protocol) connects "
              "models to tools and data. In the middle, A2A enables agent-to-agent communication — "
              "multi-turn reasoning, negotiation, and task delegation. At the top, frameworks like "
              "ADK, LangChain, and CrewAI build the actual agents. A2A is the glue between them.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Key Design Principles
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Key Design Principles", "Five foundational pillars of A2A")

principles = [
    ("Simplicity",
     BLUE2,
     ["Leverages existing web standards",
      "HTTP · JSON-RPC · SSE",
      "",
      "Minimal new concepts —",
      "familiar to any web developer"]),
    ("Enterprise\nReadiness",
     NAVY,
     ["Authentication & authorization",
      "Security & privacy",
      "Distributed tracing",
      "Monitoring & observability",
      "Production-grade from day one"]),
    ("Async\nOperations",
     TEAL,
     ["Supports long-running tasks",
      "Streaming via SSE",
      "Push notifications via webhook",
      "Human-in-the-loop workflows",
      "Non-blocking by design"]),
    ("Modality\nIndependence",
     GREEN,
     ["Not just plain text",
      "Text · Files · Structured data",
      "Bytes (raw) · URLs",
      "MIME-typed content",
      "Rich multi-modal interactions"]),
    ("Opaque\nExecution",
     RGBColor(0x6A, 0x1B, 0x9A),
     ["Agents collaborate without",
      "exposing internal logic",
      "",
      "Preserves intellectual property",
      "& implementation secrets"]),
]

px = Inches(0.35)
py = Inches(1.2)
pw = Inches(2.45)
ph = Inches(5.9)
gap = Inches(0.1)

for label, bg, lines in principles:
    add_rect(sl, px, py, pw, ph, fill=bg)
    add_textbox(sl, px + Inches(0.12), py + Inches(0.15), pw - Inches(0.24), Inches(0.65),
                label, size=16, bold=True, color=GOLD)
    y3 = py + Inches(0.9)
    for line in lines:
        add_textbox(sl, px + Inches(0.12), y3, pw - Inches(0.24), Inches(0.38),
                    line, size=12, color=WHITE)
        y3 += Inches(0.37)
    px += pw + gap

add_notes(sl, "A2A was designed around five core principles. Simplicity ensures developers can adopt "
              "it using skills they already have. Enterprise Readiness means it handles auth, security, "
              "tracing, and monitoring out of the box. Async Operations supports long-running and "
              "human-in-the-loop workflows. Modality Independence allows rich content beyond text. "
              "Opaque Execution means agents never need to reveal their internal implementation.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Core Actors
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Core Actors", "The three participants in every A2A interaction")

actors = [
    ("USER", RGBColor(0xFF, 0xF8, 0xE1), NAVY,
     "End entity initiating requests",
     ["Human operator or automated service",
      "Defines the goal or task to accomplish",
      "May be a person, a pipeline, or another system",
      "Interacts through the A2A Client"]),
    ("A2A CLIENT\n(Client Agent)", RGBColor(0xE3, 0xF2, 0xFD), NAVY,
     "Application acting on behalf of the User",
     ["AI agent, application, or service",
      "Initiates A2A communication",
      "Sends tasks to remote agents",
      "Manages authentication & context",
      "Handles streaming & polling"]),
    ("A2A SERVER\n(Remote Agent)", RGBColor(0xE8, 0xF5, 0xE9), NAVY,
     "AI agent exposing an HTTP endpoint",
     ["Implements the A2A specification",
      "Receives and processes tasks",
      "Returns messages and artifacts",
      "Opaque to the client — internal",
      "logic never exposed"]),
]

bx2 = Inches(0.35)
by2 = Inches(1.25)
bw2 = Inches(3.95)
bh2 = Inches(5.7)
gap2 = Inches(0.35)

for actor_name, bg2, fg2, subtitle2, bullets2 in actors:
    # Card
    add_rect(sl, bx2, by2, bw2, bh2, fill=bg2, line=NAVY, line_w=Pt(2))
    # Title bar
    add_rect(sl, bx2, by2, bw2, Inches(0.75), fill=NAVY)
    add_textbox(sl, bx2 + Inches(0.12), by2 + Inches(0.1), bw2 - Inches(0.24), Inches(0.55),
                actor_name, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Subtitle
    add_textbox(sl, bx2 + Inches(0.12), by2 + Inches(0.85), bw2 - Inches(0.24), Inches(0.4),
                subtitle2, size=12.5, bold=True, color=NAVY)
    # Bullets
    y4 = by2 + Inches(1.35)
    for b2 in bullets2:
        add_textbox(sl, bx2 + Inches(0.18), y4, bw2 - Inches(0.3), Inches(0.34),
                    "• " + b2, size=12, color=BLACK)
        y4 += Inches(0.36)
    bx2 += bw2 + gap2

# Arrows between actors
for ax3 in [Inches(4.38), Inches(8.73)]:
    add_rect(sl, ax3, Inches(3.85), Inches(0.27), Inches(0.35), fill=GOLD)
    add_textbox(sl, ax3 - Inches(0.03), Inches(3.83), Inches(0.35), Inches(0.38),
                "→", size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_notes(sl, "Every A2A interaction involves three actors. The User sets the goal. The A2A Client "
              "(Client Agent) acts on the user's behalf, initiating tasks and managing the communication "
              "flow. The A2A Server (Remote Agent) receives and processes tasks, returning results. "
              "The server's internal logic is completely opaque — clients only see inputs and outputs.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Fundamental Concepts
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Fundamental Concepts", "The five core building blocks of the A2A protocol")

concepts = [
    ("Agent Card",   NAVY,
     "JSON metadata document for discovery",
     ["Describes agent identity", "Lists capabilities & skills",
      "Specifies endpoint URL", "Declares auth requirements"]),
    ("Task",         BLUE2,
     "Stateful unit of work with unique ID",
     ["Facilitates long-running ops", "Supports multi-turn interactions",
      "Has defined lifecycle states", "Identified by taskId"]),
    ("Message",      TEAL,
     "Single communication turn",
     ["Role: 'user' or 'agent'", "Contains one or more Parts",
      "Conveys instructions & context", "Carries answers & requests"]),
    ("Part",         GREEN,
     "Smallest content container",
     ["TextPart — plain text", "FilePart — bytes or URL",
      "DataPart — structured JSON", "Typed by mediaType (MIME)"]),
    ("Artifact",     RGBColor(0x6A, 0x1B, 0x9A),
     "Tangible agent-generated output",
     ["Documents, images, data", "Composed of Parts",
      "Generated during task processing", "Returned as task result"]),
]

cx3 = Inches(0.35)
cy3 = Inches(1.25)
cw3 = Inches(2.45)
ch3 = Inches(5.7)
cgap = Inches(0.1)

for cname, cbg, cdesc, cbullets in concepts:
    add_rect(sl, cx3, cy3, cw3, ch3, fill=cbg)
    add_textbox(sl, cx3 + Inches(0.1), cy3 + Inches(0.12), cw3 - Inches(0.2), Inches(0.45),
                cname, size=17, bold=True, color=GOLD)
    add_rect(sl, cx3, cy3 + Inches(0.62), cw3, Inches(0.03), fill=GOLD)
    add_textbox(sl, cx3 + Inches(0.1), cy3 + Inches(0.72), cw3 - Inches(0.2), Inches(0.5),
                cdesc, size=11.5, bold=False, italic=True, color=RGBColor(0xB0, 0xC4, 0xDE))
    y5 = cy3 + Inches(1.3)
    for cb in cbullets:
        add_textbox(sl, cx3 + Inches(0.12), y5, cw3 - Inches(0.24), Inches(0.38),
                    "• " + cb, size=11.5, color=WHITE)
        y5 += Inches(0.4)
    cx3 += cw3 + cgap

add_notes(sl, "Five core concepts underpin the entire A2A protocol. The Agent Card provides discovery "
              "and capability metadata. A Task is the stateful unit of work. A Message is a single "
              "turn in a conversation. A Part is the smallest content unit inside a message. An "
              "Artifact is the tangible output produced by an agent during task processing.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Agent Card Deep Dive
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Agent Card Deep Dive", "The agent's identity & capability advertisement")

# Left description panel
add_rect(sl, Inches(0.35), Inches(1.2), Inches(5.5), Inches(5.9), fill=WHITE, line=NAVY, line_w=Pt(1.5))
add_textbox(sl, Inches(0.5), Inches(1.3), Inches(5.2), Inches(0.4),
            "What is an Agent Card?", size=15, bold=True, color=NAVY)
desc_items = [
    "A JSON metadata document that describes an agent",
    "Served at:  /.well-known/agent.json",
    "",
    "CONTAINS:",
    "  Identity — name, description, version",
    "  Capabilities — streaming, push notifications",
    "  Skills — specific tasks the agent can perform",
    "  Authentication — required schemes & scopes",
    "  Endpoint URLs — where to send requests",
    "",
    "USED FOR:",
    "  Client agent discovery (find who can help)",
    "  Capability negotiation before connecting",
    "  Auth setup before sending first message",
    "",
    "Extended Agent Card available post-auth",
    "  (reveals additional private capabilities)",
]
y6 = Inches(1.75)
for d in desc_items:
    bold6 = d.endswith(":") or d.startswith("  ")
    add_textbox(sl, Inches(0.5), y6, Inches(5.2), Inches(0.27),
                d, size=12, color=BLACK,
                bold=(d in ["What is an Agent Card?", "CONTAINS:", "USED FOR:"]))
    y6 += Inches(0.285)

# Right JSON snippet
add_rect(sl, Inches(6.1), Inches(1.2), Inches(6.9), Inches(5.9), fill=DGRAY)
add_textbox(sl, Inches(6.25), Inches(1.28), Inches(6.6), Inches(0.35),
            "// Example Agent Card (/.well-known/agent.json)",
            size=10, color=RGBColor(0x6A, 0x9F, 0x58), font_name="Courier New")

json_lines = [
    '{',
    '  "name": "TravelPlannerAgent",',
    '  "description": "Books flights and hotels",',
    '  "version": "1.0.0",',
    '  "url": "https://travel.example.com/a2a",',
    '  "capabilities": {',
    '    "streaming": true,',
    '    "pushNotifications": true',
    '  },',
    '  "skills": [',
    '    {',
    '      "id": "book-flight",',
    '      "name": "Book Flight",',
    '      "description": "Search and book flights"',
    '    }',
    '  ],',
    '  "authentication": {',
    '    "schemes": ["OAuth2"],',
    '    "credentials": null',
    '  }',
    '}',
]
y7 = Inches(1.68)
for jl in json_lines:
    col7 = RGBColor(0xD4, 0xD4, 0xD4)
    if jl.strip().startswith('"') and ':' in jl:
        parts7 = jl.split(':', 1)
        col7 = RGBColor(0x9C, 0xDC, 0xFE)
    add_textbox(sl, Inches(6.25), y7, Inches(6.6), Inches(0.27),
                jl, size=10.5, color=col7, font_name="Courier New")
    y7 += Inches(0.265)

add_notes(sl, "The Agent Card is the discovery mechanism at the heart of A2A. It lives at a "
              "well-known URL and describes everything a client needs to know before connecting: "
              "who the agent is, what it can do, how to authenticate, and where to send requests. "
              "An extended card with additional capabilities is available after authentication.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Task Lifecycle (grid-aligned)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "The Task Lifecycle", "States and transitions in A2A task processing")

# ── Grid constants ────────────────────────────────────────────────────────────
# Slide: 13.33" × 7.5" | Header: ~1.1" | Content: 1.2" – 7.2"
BOX_W  = Inches(2.15)
BOX_H  = Inches(0.88)

# Column left-edges (3 columns)
CX_A = Inches(0.40)   # submitted / input-required
CX_B = Inches(3.50)   # working
CX_C = Inches(9.10)   # completed / failed / canceled

# Row top-edges (3 rows)
CY_TOP = Inches(1.55)   # completed
CY_MID = Inches(3.05)   # submitted / working / failed  ← main flow row
CY_BOT = Inches(4.55)   # canceled
CY_INP = Inches(4.55)   # input-required (same row as canceled, under working)

# Vertical centers (for connector alignment)
VC_TOP = CY_TOP + BOX_H / 2   # 1.99"
VC_MID = CY_MID + BOX_H / 2   # 3.49"
VC_BOT = CY_BOT + BOX_H / 2   # 4.99"
VC_INP = CY_INP + BOX_H / 2   # 4.99"  (same as VC_BOT, different column)

# Right-edges (for connector start/end)
RE_A = CX_A + BOX_W   # 2.55"
RE_B = CX_B + BOX_W   # 5.65"
LE_C = CX_C            # 9.10"

# Vertical spine x (between Col B and Col C)
SPINE_X = Inches(7.80)

# ── Helper: draw one state box ────────────────────────────────────────────────
def state_box(slide, x, y, label, bg, fg, sub):
    add_rect(slide, x, y, BOX_W, BOX_H, fill=bg)
    add_textbox(slide, x + Inches(0.08), y + Inches(0.07),
                BOX_W - Inches(0.16), Inches(0.42),
                label, size=13, bold=True, color=fg, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.05), y + Inches(0.50),
                BOX_W - Inches(0.10), Inches(0.33),
                sub, size=9, color=fg, align=PP_ALIGN.CENTER)

# ── Helper: thin rectangle as a line segment ──────────────────────────────────
LINE_CLR = RGBColor(0x88, 0x88, 0x88)
LINE_T   = Inches(0.04)   # line thickness

def hline(slide, x, y, length, color=LINE_CLR):
    """Horizontal line centred on y."""
    add_rect(slide, x, y - LINE_T / 2, length, LINE_T, fill=color)

def vline(slide, x, y, height, color=LINE_CLR):
    """Vertical line centred on x."""
    add_rect(slide, x - LINE_T / 2, y, LINE_T, height, fill=color)

def arrow_r(slide, x, y, color=LINE_CLR):
    """Right-pointing arrowhead at (x, y)."""
    add_textbox(slide, x - Inches(0.08), y - Inches(0.18),
                Inches(0.25), Inches(0.36), "▶", size=11, color=color,
                align=PP_ALIGN.LEFT)

def arrow_d(slide, x, y, color=LINE_CLR):
    """Down-pointing arrowhead at (x, y)."""
    add_textbox(slide, x - Inches(0.14), y - Inches(0.05),
                Inches(0.28), Inches(0.30), "▼", size=11, color=color,
                align=PP_ALIGN.CENTER)

# ── Draw the six state boxes ──────────────────────────────────────────────────
state_box(sl, CX_A, CY_MID, "submitted",
          NAVY,                      WHITE, "received · not yet started")
state_box(sl, CX_B, CY_MID, "working",
          BLUE2,                     WHITE, "agent actively processing")
state_box(sl, CX_C, CY_TOP, "completed",
          GREEN,                     WHITE, "success · artifacts ready")
state_box(sl, CX_C, CY_MID, "failed",
          RGBColor(0xC0, 0x27, 0x27), WHITE, "terminal error state")
state_box(sl, CX_C, CY_BOT, "canceled",
          MGRAY,                     WHITE, "canceled by client")
state_box(sl, CX_B, CY_INP, "input-required",
          GOLD,                      NAVY,  "needs more info · multi-turn")

# ── Arrow 1: submitted → working (horizontal, main flow row) ─────────────────
hline(sl,  RE_A, VC_MID,  CX_B - RE_A)
arrow_r(sl, CX_B, VC_MID)

# ── Arrow 2: working → vertical spine (horizontal stub) ──────────────────────
hline(sl, RE_B, VC_MID, SPINE_X - RE_B)

# ── Vertical spine: from VC_TOP to VC_BOT ────────────────────────────────────
vline(sl, SPINE_X, VC_TOP, VC_BOT - VC_TOP)

# ── Arrow 3: spine → completed ────────────────────────────────────────────────
hline(sl, SPINE_X, VC_TOP, LE_C - SPINE_X)
arrow_r(sl, LE_C, VC_TOP)

# ── Arrow 4: spine → failed (same y as main flow) ────────────────────────────
hline(sl, SPINE_X, VC_MID, LE_C - SPINE_X)
arrow_r(sl, LE_C, VC_MID)

# ── Arrow 5: spine → canceled ─────────────────────────────────────────────────
hline(sl, SPINE_X, VC_BOT, LE_C - SPINE_X)
arrow_r(sl, LE_C, VC_BOT)

# ── Arrow 6: working ↓ input-required (vertical, under Col B) ─────────────────
WRK_CX = CX_B + BOX_W / 2   # centre of working / input-required boxes
vline(sl, WRK_CX, CY_MID + BOX_H, CY_INP - (CY_MID + BOX_H), color=GOLD)
arrow_d(sl, WRK_CX, CY_INP, color=GOLD)

# ── "resumes" label (input-required loops back to working) ────────────────────
add_textbox(sl, CX_A + Inches(0.05), CY_INP + Inches(0.18),
            Inches(2.85), Inches(0.30),
            "◀ resumes working", size=9, bold=True, color=GOLD,
            align=PP_ALIGN.CENTER)

# ── State legend (2-column, 3 rows) at bottom ────────────────────────────────
LG_Y    = Inches(5.95)
LG_H    = Inches(1.08)
LG_LBL  = Inches(1.55)   # label column width
LG_TXT  = Inches(4.65)   # description column width
LG_GAP  = Inches(0.20)   # gap between columns
COL1_X  = Inches(0.40)
COL2_X  = COL1_X + LG_LBL + LG_TXT + LG_GAP + Inches(0.30)
ROW_H   = Inches(0.30)

add_rect(sl, Inches(0.35), LG_Y, Inches(12.63), LG_H,
         fill=RGBColor(0xF4, 0xF6, 0xF9), line=NAVY, line_w=Pt(0.75))

legend_items = [
    # (label, color, description text)
    ("submitted",      NAVY,                       "Task received by server; awaiting processing"),
    ("working",        BLUE2,                      "Agent actively processes; may stream incremental updates"),
    ("input-required", GOLD,                       "Agent paused; client sends follow-up to resume the task"),
    ("completed",      GREEN,                      "Finished successfully; artifacts & messages available"),
    ("failed",         RGBColor(0xC0, 0x27, 0x27),"Terminal error; check task response for details"),
    ("canceled",       MGRAY,                      "Client called Cancel Task; terminal state"),
]

for i, (lbl, clr, desc) in enumerate(legend_items):
    col_x = COL1_X if i % 2 == 0 else COL2_X
    row_y = LG_Y + Inches(0.12) + (i // 2) * ROW_H
    add_textbox(sl, col_x, row_y, LG_LBL, ROW_H - Inches(0.04),
                f"● {lbl}:", size=9.5, bold=True, color=clr)
    add_textbox(sl, col_x + LG_LBL, row_y, LG_TXT, ROW_H - Inches(0.04),
                desc, size=9.5, color=BLACK)

add_notes(sl, "A task in A2A goes through a well-defined lifecycle. It starts as 'submitted', moves "
              "to 'working' when the agent begins processing, and then reaches a terminal state: "
              "'completed' (success), 'failed' (error), or 'canceled' (client-initiated). The "
              "'input-required' state enables multi-turn interactions where the agent needs more "
              "information before it can continue.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Messages & Parts
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Messages & Parts", "The structure of communication in A2A")

# Message structure diagram (top half)
add_rect(sl, Inches(0.35), Inches(1.2), Inches(12.6), Inches(1.7), fill=WHITE, line=NAVY, line_w=Pt(1.5))
add_textbox(sl, Inches(0.5), Inches(1.3), Inches(5), Inches(0.35),
            "MESSAGE STRUCTURE", size=13, bold=True, color=NAVY)

msg_fields = [
    ("messageId", "Unique identifier for the message", Inches(0.5)),
    ("taskId",    "Links message to its parent task",   Inches(3.8)),
    ("role",      "\"user\" or \"agent\"",               Inches(7.1)),
    ("parts",     "List of Part objects (content)",     Inches(10.0)),
]
for fname, fdesc, fx9 in msg_fields:
    add_rect(sl, fx9, Inches(1.7), Inches(2.8), Inches(0.65), fill=LGRAY, line=NAVY, line_w=Pt(0.75))
    add_textbox(sl, fx9 + Inches(0.08), Inches(1.73), Inches(2.65), Inches(0.28),
                fname, size=12, bold=True, color=NAVY, font_name="Courier New")
    add_textbox(sl, fx9 + Inches(0.08), Inches(2.0), Inches(2.65), Inches(0.3),
                fdesc, size=10.5, color=MGRAY)

# Part types (bottom half)
part_types = [
    ("TextPart",   NAVY,  WHITE,
     "Plain text content",
     ["{ \"kind\": \"text\",", "  \"text\": \"Book a flight to Paris\" }"],
     "mediaType: text/plain"),
    ("FilePart",   TEAL,  WHITE,
     "File by URL or raw bytes",
     ["{ \"kind\": \"file\",", "  \"file\": {",
      "    \"url\": \"https://...\",", "    \"mimeType\": \"image/png\" } }"],
     "mediaType: image/*, application/pdf, etc."),
    ("DataPart",   BLUE2, WHITE,
     "Structured JSON data",
     ["{ \"kind\": \"data\",", "  \"data\": {",
      "    \"booking_ref\": \"XY1234\",", "    \"price\": 499.00 } }"],
     "mediaType: application/json"),
    ("Artifact",   GREEN, WHITE,
     "Agent-generated output\nComposed of Parts",
     ["{ \"artifactId\": \"doc-1\",", "  \"name\": \"Itinerary\",",
      "  \"parts\": [ ... ] }"],
     "Terminal task output"),
]

px9 = Inches(0.35)
py9 = Inches(3.1)
pw9 = Inches(3.05)
ph9 = Inches(4.0)
gap9 = Inches(0.12)

for ptype, pbg, pfg, pdesc, pcode, pnote in part_types:
    add_rect(sl, px9, py9, pw9, ph9, fill=pbg)
    add_textbox(sl, px9 + Inches(0.12), py9 + Inches(0.1), pw9 - Inches(0.24), Inches(0.38),
                ptype, size=15, bold=True, color=GOLD)
    add_textbox(sl, px9 + Inches(0.12), py9 + Inches(0.52), pw9 - Inches(0.24), Inches(0.4),
                pdesc, size=11.5, italic=True, color=RGBColor(0xB0, 0xC4, 0xDE))
    # Code snippet
    yc = py9 + Inches(1.0)
    for cl in pcode:
        add_textbox(sl, px9 + Inches(0.12), yc, pw9 - Inches(0.24), Inches(0.26),
                    cl, size=10, color=RGBColor(0xD4, 0xD4, 0xD4), font_name="Courier New")
        yc += Inches(0.27)
    add_textbox(sl, px9 + Inches(0.12), py9 + ph9 - Inches(0.4), pw9 - Inches(0.24), Inches(0.35),
                pnote, size=10.5, italic=True, color=GOLD)
    px9 += pw9 + gap9

add_notes(sl, "A Message is a single turn in a conversation, identified by messageId, linked to a "
              "Task, and tagged with a role (user or agent). It contains one or more Parts. "
              "TextPart carries plain text. FilePart carries files by URL or raw bytes. DataPart "
              "carries structured JSON. Artifacts are agent-generated outputs composed of Parts — "
              "they represent the final deliverable of a task.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Interaction Patterns
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Interaction Patterns", "Three ways clients receive updates from servers")

patterns = [
    ("Request / Response\n(Polling)",
     NAVY, WHITE, GOLD,
     "Client sends request; server responds synchronously. "
     "For long-running tasks, client polls periodically using Get Task.",
     ["Simple to implement", "Works with any HTTP client",
      "Higher latency for long tasks", "Periodic Get Task calls",
      "Good for infrequent updates"],
     ["1. Client → sendMessage", "2. Server returns Task (working)",
      "3. Client polls: getTask(taskId)", "4. Repeat until terminal state",
      "5. Retrieve artifacts"]),
    ("Streaming\n(Server-Sent Events)",
     TEAL, WHITE, GOLD,
     "Server streams real-time incremental results over a persistent "
     "HTTP connection using Server-Sent Events (SSE).",
     ["Real-time updates", "Persistent HTTP connection",
      "Lower latency", "Ideal for long responses",
      "Automatic reconnect support"],
     ["1. Client → sendMessageStream", "2. Connection stays open",
      "3. Server pushes SSE events", "4. Events: task updates + artifacts",
      "5. Stream closes on completion"]),
    ("Push Notifications\n(Webhooks)",
     GREEN, WHITE, GOLD,
     "Server sends asynchronous updates to a client-provided webhook "
     "URL. Ideal for extended operations where the client cannot "
     "maintain a connection.",
     ["Fully asynchronous", "No persistent connection needed",
      "Client-provided webhook URL", "Best for very long tasks",
      "Human-in-the-loop workflows"],
     ["1. Client registers webhook URL", "2. Client → sendMessage",
      "3. Server sends task updates", "   via HTTP POST to webhook",
      "4. Client processes async events"]),
]

px10 = Inches(0.35)
py10 = Inches(1.2)
pw10 = Inches(4.1)
ph10 = Inches(6.0)
gap10 = Inches(0.2)

for pname, pbg, pfg, pacc, pdesc, ppros, psteps in patterns:
    add_rect(sl, px10, py10, pw10, ph10, fill=pbg)
    # Title
    add_textbox(sl, px10 + Inches(0.12), py10 + Inches(0.1), pw10 - Inches(0.24), Inches(0.6),
                pname, size=16, bold=True, color=pacc)
    # Divider
    add_rect(sl, px10, py10 + Inches(0.75), pw10, Inches(0.025), fill=pacc)
    # Description
    add_textbox(sl, px10 + Inches(0.12), py10 + Inches(0.82), pw10 - Inches(0.24), Inches(0.8),
                pdesc, size=11, italic=True, color=RGBColor(0xB0, 0xC4, 0xDE))
    # Pros header
    add_textbox(sl, px10 + Inches(0.12), py10 + Inches(1.7), pw10 - Inches(0.24), Inches(0.3),
                "Characteristics:", size=11.5, bold=True, color=pacc)
    y_p = py10 + Inches(2.05)
    for pro in ppros:
        add_textbox(sl, px10 + Inches(0.18), y_p, pw10 - Inches(0.3), Inches(0.27),
                    "• " + pro, size=11, color=WHITE)
        y_p += Inches(0.3)
    # Steps header
    add_textbox(sl, px10 + Inches(0.12), py10 + Inches(3.62), pw10 - Inches(0.24), Inches(0.3),
                "Flow:", size=11.5, bold=True, color=pacc)
    y_s = py10 + Inches(3.97)
    for step in psteps:
        add_textbox(sl, px10 + Inches(0.18), y_s, pw10 - Inches(0.3), Inches(0.27),
                    step, size=10.5, color=RGBColor(0xD4, 0xD4, 0xD4))
        y_s += Inches(0.3)
    px10 += pw10 + gap10

add_notes(sl, "A2A supports three interaction patterns. Request/Response with Polling is the simplest: "
              "send a message, poll for updates. Streaming via Server-Sent Events delivers real-time "
              "incremental updates over a persistent connection — ideal for long responses. Push "
              "Notifications use webhooks for fully asynchronous delivery, best for very long "
              "operations or human-in-the-loop workflows where maintaining a connection isn't practical.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — The 11 Core Operations
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "The 11 Core Operations", "Complete A2A specification API surface")

operations = [
    ("Messaging",    NAVY,
     [("1", "Send Message",         "Primary operation; returns Task or direct Message"),
      ("2", "Send Streaming Msg",   "Like Send Message with real-time SSE streaming")]),
    ("Task Mgmt",    BLUE2,
     [("3", "Get Task",             "Retrieves current state and details of a task"),
      ("4", "List Tasks",           "Retrieves task list with optional filtering & pagination"),
      ("5", "Cancel Task",          "Requests cancellation of an ongoing task"),
      ("6", "Subscribe to Task",    "Establishes streaming connection for task update events")]),
    ("Push Notif.",  TEAL,
     [("7",  "Create Push Config",  "Sets up a webhook URL for async task updates"),
      ("8",  "Get Push Config",     "Retrieves an existing webhook configuration"),
      ("9",  "List Push Configs",   "Lists all webhook configs for a specific task"),
      ("10", "Delete Push Config",  "Removes a webhook configuration")]),
    ("Discovery",    GREEN,
     [("11", "Get Extended Card",   "Retrieves detailed agent capabilities post-authentication")]),
]

# Header row
hx = Inches(0.35)
hy = Inches(1.25)
col_widths = [Inches(1.4), Inches(0.45), Inches(3.0), Inches(7.4)]
headers = ["Category", "#", "Operation", "Description"]
hbgs = [NAVY, NAVY, NAVY, NAVY]
for hdr, hw, hbg in zip(headers, col_widths, hbgs):
    add_rect(sl, hx, hy, hw, Inches(0.38), fill=hbg)
    add_textbox(sl, hx + Inches(0.05), hy + Inches(0.04), hw - Inches(0.1), Inches(0.3),
                hdr, size=12, bold=True, color=WHITE)
    hx += hw

# Data rows
row_y = hy + Inches(0.38)
row_h = Inches(0.48)
alt_bg = RGBColor(0xE8, 0xEE, 0xF8)

for cat_name, cat_color, ops_list in operations:
    # Category spanning cell
    cat_h = row_h * len(ops_list)
    add_rect(sl, Inches(0.35), row_y, Inches(1.4), cat_h, fill=cat_color)
    add_textbox(sl, Inches(0.38), row_y + cat_h/2 - Inches(0.2),
                Inches(1.35), Inches(0.4),
                cat_name, size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, (num, op_name, op_desc) in enumerate(ops_list):
        bg_row = WHITE if i % 2 == 0 else alt_bg
        # Number
        add_rect(sl, Inches(1.75), row_y, Inches(0.45), row_h, fill=bg_row, line=LGRAY, line_w=Pt(0.5))
        add_textbox(sl, Inches(1.78), row_y + Inches(0.1), Inches(0.38), Inches(0.28),
                    num, size=12, bold=True, color=cat_color, align=PP_ALIGN.CENTER)
        # Op name
        add_rect(sl, Inches(2.2), row_y, Inches(3.0), row_h, fill=bg_row, line=LGRAY, line_w=Pt(0.5))
        add_textbox(sl, Inches(2.28), row_y + Inches(0.1), Inches(2.9), Inches(0.28),
                    op_name, size=12, bold=True, color=cat_color)
        # Description
        add_rect(sl, Inches(5.2), row_y, Inches(7.4), row_h, fill=bg_row, line=LGRAY, line_w=Pt(0.5))
        add_textbox(sl, Inches(5.28), row_y + Inches(0.1), Inches(7.25), Inches(0.28),
                    op_desc, size=11.5, color=BLACK)
        row_y += row_h

# Idempotency note
add_rect(sl, Inches(0.35), row_y + Inches(0.1), Inches(12.6), Inches(0.38), fill=RGBColor(0xFF,0xF8,0xE1), line=GOLD, line_w=Pt(1))
add_textbox(sl, Inches(0.5), row_y + Inches(0.15), Inches(12.3), Inches(0.28),
            "Note: Get operations and Cancel Task are idempotent. Send Message MAY be idempotent via messageId. "
            "Blocking mode supported via SendMessageConfiguration.",
            size=10.5, color=NAVY)

add_notes(sl, "A2A defines exactly 11 operations across four categories. The two Messaging operations "
              "are the primary workhorses. The four Task Management operations handle the full task "
              "lifecycle. The four Push Notification operations manage webhooks. The single Discovery "
              "operation retrieves extended agent capabilities after authentication.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Protocol Bindings
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Protocol Bindings", "Three transport options for A2A communication")

bindings = [
    ("JSON-RPC 2.0",   NAVY,
     "Method-based RPC over HTTP",
     ["Standard JSON-RPC 2.0 envelope",
      "Method names map to operations",
      "Request/Response and Notification",
      "Easy to debug with standard tools",
      "Widely supported in all languages"],
     ["Best for:", "• Most general use cases",
      "• Web and cloud environments",
      "• Teams familiar with REST/JSON",
      "• Tooling & debugging ease"],
     ['{"jsonrpc": "2.0",',
      ' "method": "tasks/send",',
      ' "params": { ... },',
      ' "id": "req-001"}']),
    ("gRPC",           TEAL,
     "High-performance binary RPC",
     ["Protocol Buffers (binary encoding)",
      "HTTP/2 multiplexing",
      "Strongly typed contracts",
      "Native streaming support",
      "Auto-generated client stubs"],
     ["Best for:", "• High-throughput scenarios",
      "• Low-latency requirements",
      "• Microservice architectures",
      "• Multi-language polyglot systems"],
     ["service A2AService {",
      "  rpc SendMessage(",
      "    MessageRequest)",
      "    returns (Task);",
      "}"]),
    ("HTTP+JSON/REST",  GREEN,
     "Standard REST API with JSON",
     ["Familiar HTTP verbs (POST, GET)",
      "JSON request/response bodies",
      "OpenAPI/Swagger compatible",
      "Stateless by default",
      "Universal tooling support"],
     ["Best for:", "• Developer familiarity",
      "• Integration with existing APIs",
      "• Simple client implementations",
      "• API gateway compatibility"],
     ["POST /tasks/send",
      "Content-Type: application/json",
      "",
      '{ "message": {',
      '    "role": "user",',
      '    "parts": [...] } }']),
]

bx11 = Inches(0.35)
by11 = Inches(1.25)
bw11 = Inches(4.1)
bh11 = Inches(5.95)
bgap = Inches(0.2)

for bname, bbg, bsubhead, bfeats, bbest, bcode in bindings:
    add_rect(sl, bx11, by11, bw11, bh11, fill=bbg)
    add_textbox(sl, bx11 + Inches(0.12), by11 + Inches(0.1), bw11 - Inches(0.24), Inches(0.45),
                bname, size=18, bold=True, color=GOLD)
    add_textbox(sl, bx11 + Inches(0.12), by11 + Inches(0.58), bw11 - Inches(0.24), Inches(0.32),
                bsubhead, size=12, italic=True, color=RGBColor(0xB0, 0xC4, 0xDE))
    add_rect(sl, bx11, by11 + Inches(0.95), bw11, Inches(0.025), fill=GOLD)
    # Features
    yf = by11 + Inches(1.05)
    for feat in bfeats:
        add_textbox(sl, bx11 + Inches(0.18), yf, bw11 - Inches(0.3), Inches(0.27),
                    "• " + feat, size=11, color=WHITE)
        yf += Inches(0.28)
    # Best for
    add_textbox(sl, bx11 + Inches(0.12), by11 + Inches(2.55), bw11 - Inches(0.24), Inches(0.3),
                bbest[0], size=11.5, bold=True, color=GOLD)
    yb = by11 + Inches(2.87)
    for bl in bbest[1:]:
        add_textbox(sl, bx11 + Inches(0.18), yb, bw11 - Inches(0.3), Inches(0.27),
                    bl, size=11, color=WHITE)
        yb += Inches(0.28)
    # Code
    add_rect(sl, bx11 + Inches(0.1), by11 + Inches(4.15), bw11 - Inches(0.2), Inches(1.6), fill=DGRAY)
    ycc = by11 + Inches(4.22)
    for cl in bcode:
        add_textbox(sl, bx11 + Inches(0.18), ycc, bw11 - Inches(0.35), Inches(0.27),
                    cl, size=9.5, color=RGBColor(0xD4, 0xD4, 0xD4), font_name="Courier New")
        ycc += Inches(0.26)
    bx11 += bw11 + bgap

add_notes(sl, "A2A supports three protocol bindings. JSON-RPC 2.0 is the primary binding used by "
              "most SDKs — it maps naturally to A2A operations and is easy to debug. gRPC provides "
              "high performance with binary encoding and native streaming, ideal for "
              "high-throughput microservices. HTTP+JSON/REST is the most familiar option, compatible "
              "with standard web tooling and API gateways.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Authentication & Security
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Authentication & Security", "Enterprise-grade security built into the protocol")

auth_schemes = [
    ("API Key",       NAVY,
     "Simple secret key passed in HTTP headers",
     ["X-API-Key: <key> header", "Fast to implement",
      "Best for server-to-server", "Rotate regularly"]),
    ("HTTP Basic /\nBearer Token", BLUE2,
     "Standard HTTP authentication schemes",
     ["Authorization: Bearer <token>", "JWT or opaque tokens",
      "Widely supported", "Use with TLS always"]),
    ("OAuth 2.0",     TEAL,
     "Full OAuth 2.0 with multiple grant types",
     ["Authorization Code flow", "Client Credentials flow",
      "Device Code flow", "Short-lived access tokens"]),
    ("OpenID Connect", GREEN,
     "OAuth 2.0 + identity layer (JWT)",
     ["ID token with user claims", "Discovery via .well-known",
      "Used in A2A request lifecycle", "Standard JWT verification"]),
    ("Mutual TLS\n(mTLS)", RGBColor(0x6A, 0x1B, 0x9A),
     "Certificate-based bidirectional auth",
     ["Both client & server present certs", "Strongest auth guarantee",
      "Zero-trust network model", "Enterprise & regulated industries"]),
]

ax12 = Inches(0.35)
ay12 = Inches(1.25)
aw12 = Inches(2.42)
ah12 = Inches(4.2)
agap = Inches(0.09)

for aname, abg, adesc, abullets in auth_schemes:
    add_rect(sl, ax12, ay12, aw12, ah12, fill=abg)
    add_textbox(sl, ax12 + Inches(0.1), ay12 + Inches(0.1), aw12 - Inches(0.2), Inches(0.55),
                aname, size=13.5, bold=True, color=GOLD)
    add_rect(sl, ax12, ay12 + Inches(0.7), aw12, Inches(0.025), fill=GOLD)
    add_textbox(sl, ax12 + Inches(0.1), ay12 + Inches(0.78), aw12 - Inches(0.2), Inches(0.55),
                adesc, size=10.5, italic=True, color=RGBColor(0xB0, 0xC4, 0xDE))
    yb12 = ay12 + Inches(1.42)
    for ab in abullets:
        add_textbox(sl, ax12 + Inches(0.12), yb12, aw12 - Inches(0.24), Inches(0.3),
                    "• " + ab, size=10.5, color=WHITE)
        yb12 += Inches(0.33)
    ax12 += aw12 + agap

# How it works strip
add_rect(sl, Inches(0.35), Inches(5.55), Inches(12.6), Inches(1.55), fill=WHITE, line=NAVY, line_w=Pt(1.5))
add_textbox(sl, Inches(0.5), Inches(5.63), Inches(5), Inches(0.35),
            "How Authentication Works in A2A", size=13, bold=True, color=NAVY)
how_items = [
    ("Step 1:", "Client fetches Agent Card — auth schemes declared in card"),
    ("Step 2:", "Client obtains credentials (e.g., JWT via OpenID Connect)"),
    ("Step 3:", "Credentials passed via HTTP headers on every request"),
    ("Step 4:", "Server validates credentials before processing the task"),
]
hx12 = Inches(0.5)
hy12 = Inches(6.0)
for hlabel, hdesc in how_items:
    add_textbox(sl, hx12, hy12, Inches(0.75), Inches(0.28), hlabel, size=10.5, bold=True, color=NAVY)
    add_textbox(sl, hx12 + Inches(0.7), hy12, Inches(2.5), Inches(0.28), hdesc, size=10.5, color=BLACK)
    hx12 += Inches(3.22)

add_notes(sl, "A2A supports five authentication schemes to cover all enterprise use cases. API Key "
              "is the simplest. Bearer tokens cover most web use cases. OAuth 2.0 provides full "
              "delegated authorization with multiple grant types. OpenID Connect adds identity. "
              "Mutual TLS provides the strongest guarantee for regulated industries. "
              "Auth schemes are declared in the Agent Card and credentials are passed via HTTP headers.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Multi-Turn Interactions
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Multi-Turn Interactions", "Maintaining context across a conversation")

# Left — concepts
add_rect(sl, Inches(0.35), Inches(1.2), Inches(5.8), Inches(5.9), fill=WHITE, line=NAVY, line_w=Pt(1.5))
concepts_mt = [
    ("contextId", NAVY,
     "Server-generated ID grouping related Tasks and Messages logically. "
     "Enables the agent to maintain conversation state across multiple task turns."),
    ("taskId", BLUE2,
     "Unique identifier for each stateful unit of work. "
     "Allows polling, streaming subscriptions, and cancellations."),
    ("referenceTaskIds", TEAL,
     "Messages can reference related tasks by ID. "
     "Enables chaining and dependency tracking between agent tasks."),
    ("input-required", GREEN,
     "Special task state where the agent pauses and requests additional input. "
     "Client sends a follow-up message to resume the task."),
]
yc16 = Inches(1.3)
for cname16, ccol, cdesc16 in concepts_mt:
    add_textbox(sl, Inches(0.5), yc16, Inches(5.5), Inches(0.33),
                cname16, size=14, bold=True, color=ccol, font_name="Courier New")
    add_textbox(sl, Inches(0.5), yc16 + Inches(0.35), Inches(5.5), Inches(0.7),
                cdesc16, size=11.5, color=BLACK)
    yc16 += Inches(1.3)

# Right — conversation flow diagram
add_rect(sl, Inches(6.4), Inches(1.2), Inches(6.6), Inches(5.9), fill=LGRAY, line=NAVY, line_w=Pt(1.5))
add_textbox(sl, Inches(6.6), Inches(1.3), Inches(6.2), Inches(0.35),
            "Multi-Turn Conversation Flow", size=13, bold=True, color=NAVY)

flow_steps = [
    (BLUE2,  "Client",  "→", NAVY,  "Server",
     "sendMessage(message='Plan a Paris trip')", "Task created: working"),
    (NAVY,   "Server",  "→", BLUE2, "Client",
     "TaskUpdate: input-required", "'What dates? Budget?'"),
    (BLUE2,  "Client",  "→", NAVY,  "Server",
     "sendMessage(message='July 10-17, $2000')", "Same contextId"),
    (NAVY,   "Server",  "→", BLUE2, "Client",
     "TaskUpdate: working → completed", "Artifact: full itinerary"),
]

yf16 = Inches(1.8)
step_h = Inches(1.05)
for src_col, src, arrow, dst_col, dst, msg1, msg2 in flow_steps:
    # From box
    add_rect(sl, Inches(6.55), yf16, Inches(1.3), Inches(0.38), fill=src_col)
    add_textbox(sl, Inches(6.58), yf16 + Inches(0.05), Inches(1.25), Inches(0.28),
                src, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow
    add_textbox(sl, Inches(7.9), yf16 + Inches(0.05), Inches(0.5), Inches(0.28),
                arrow, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # To box
    add_rect(sl, Inches(8.45), yf16, Inches(1.3), Inches(0.38), fill=dst_col)
    add_textbox(sl, Inches(8.48), yf16 + Inches(0.05), Inches(1.25), Inches(0.28),
                dst, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Message
    add_textbox(sl, Inches(6.55), yf16 + Inches(0.42), Inches(6.25), Inches(0.28),
                msg1, size=10.5, italic=True, color=NAVY)
    add_textbox(sl, Inches(6.55), yf16 + Inches(0.7), Inches(6.25), Inches(0.28),
                msg2, size=10, color=MGRAY)
    yf16 += step_h

# contextId label
add_rect(sl, Inches(12.3), Inches(1.8), Inches(0.5), Inches(4.2), fill=GOLD)
add_textbox(sl, Inches(12.32), Inches(3.3), Inches(0.45), Inches(1.2),
            "same contextId", size=8, bold=True, color=NAVY)

add_notes(sl, "Multi-turn interactions are central to A2A. The contextId groups all related tasks "
              "and messages, giving the server the context it needs to maintain conversation state. "
              "The input-required task state enables the agent to pause and ask for more information. "
              "The referenceTaskIds field lets messages explicitly reference related work, enabling "
              "complex chained agent workflows.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Request Lifecycle
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Request Lifecycle", "Four stages from discovery to streaming results")

stages = [
    ("1", "Agent\nDiscovery", NAVY,
     ["Client fetches Agent Card",
      "from /.well-known/agent.json",
      "",
      "Learns: capabilities,",
      "endpoint URL, auth schemes,",
      "supported skills"],
     "GET /.well-known/agent.json\n→ Returns: AgentCard JSON"),
    ("2", "Authentication", BLUE2,
     ["Client inspects auth schemes",
      "declared in Agent Card",
      "",
      "Obtains JWT token via",
      "OpenID Connect (if required)",
      "or uses API Key / Bearer"],
     "POST /token\n→ Returns: JWT access_token"),
    ("3", "sendMessage\nAPI", TEAL,
     ["Client sends synchronous",
      "request with Task payload",
      "",
      "Server returns Task object",
      "with current state",
      "(working, completed, etc.)"],
     "POST /a2a\n{\"method\":\"tasks/send\"}\n→ Returns: Task"),
    ("4", "sendMessage\nStream API", GREEN,
     ["Client opens streaming",
      "connection for real-time",
      "updates via SSE",
      "",
      "Server pushes TaskUpdate",
      "events until completion"],
     "POST /a2a\n{\"method\":\"tasks/sendSubscribe\"}\n→ SSE stream of events"),
]

sx17 = Inches(0.35)
sy17 = Inches(1.25)
sw17 = Inches(3.05)
sh17 = Inches(5.9)
sgap = Inches(0.12)

for snum, stitle, sbg, sitems, scode in stages:
    add_rect(sl, sx17, sy17, sw17, sh17, fill=sbg)
    # Number badge
    add_rect(sl, sx17 + Inches(0.12), sy17 + Inches(0.12), Inches(0.45), Inches(0.45), fill=GOLD)
    add_textbox(sl, sx17 + Inches(0.12), sy17 + Inches(0.12), Inches(0.45), Inches(0.45),
                snum, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(sl, sx17 + Inches(0.65), sy17 + Inches(0.15), sw17 - Inches(0.8), Inches(0.52),
                stitle, size=15, bold=True, color=GOLD)
    add_rect(sl, sx17, sy17 + Inches(0.75), sw17, Inches(0.025), fill=GOLD)
    yp17 = sy17 + Inches(0.85)
    for item17 in sitems:
        add_textbox(sl, sx17 + Inches(0.15), yp17, sw17 - Inches(0.3), Inches(0.33),
                    item17, size=12, color=WHITE)
        yp17 += Inches(0.33)
    # Code box at bottom
    add_rect(sl, sx17 + Inches(0.1), sy17 + sh17 - Inches(1.2), sw17 - Inches(0.2), Inches(1.1), fill=DGRAY)
    y_c17 = sy17 + sh17 - Inches(1.15)
    for cl17 in scode.split('\n'):
        add_textbox(sl, sx17 + Inches(0.18), y_c17, sw17 - Inches(0.35), Inches(0.28),
                    cl17, size=9.5, color=RGBColor(0xD4, 0xD4, 0xD4), font_name="Courier New")
        y_c17 += Inches(0.27)
    sx17 += sw17 + sgap

add_notes(sl, "Every A2A interaction follows a four-stage lifecycle. Stage 1 is discovery: the "
              "client fetches the Agent Card to learn about the agent's capabilities and auth "
              "requirements. Stage 2 is authentication: the client obtains credentials. Stage 3 is "
              "the sendMessage API: a synchronous request that creates or continues a task. Stage 4 "
              "is the streaming API: an SSE connection for real-time task updates.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — SDK & Framework Support
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "SDK & Framework Support", "Official SDKs and compatible agent frameworks")

# Official SDKs section
add_rect(sl, Inches(0.35), Inches(1.2), Inches(12.6), Inches(0.42), fill=NAVY)
add_textbox(sl, Inches(0.5), Inches(1.27), Inches(12), Inches(0.28),
            "Official SDKs  (open source — github.com/google-a2a)", size=13, bold=True, color=WHITE)

sdk_items = [
    ("Python",   "github.com/google-a2a/a2a-python",   BLUE2),
    ("JavaScript", "github.com/google-a2a/a2a-js",     RGBColor(0xF5, 0xA6, 0x23)),
    ("Java",     "github.com/google-a2a/a2a-java",      RGBColor(0xC6, 0x28, 0x28)),
    ("C# / .NET", "github.com/google-a2a/a2a-dotnet",  RGBColor(0x51, 0x2B, 0xD4)),
    ("Golang",   "github.com/google-a2a/a2a-go",        TEAL),
]

sx18 = Inches(0.35)
sy18 = Inches(1.65)
sw18 = Inches(2.42)
sh18 = Inches(1.2)
sgap18 = Inches(0.08)

for sdk_name, sdk_url, sdk_color in sdk_items:
    add_rect(sl, sx18, sy18, sw18, sh18, fill=sdk_color)
    add_textbox(sl, sx18 + Inches(0.12), sy18 + Inches(0.12), sw18 - Inches(0.24), Inches(0.42),
                sdk_name, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, sx18 + Inches(0.08), sy18 + Inches(0.58), sw18 - Inches(0.16), Inches(0.5),
                sdk_url, size=8.5, color=RGBColor(0xD4, 0xD4, 0xD4),
                align=PP_ALIGN.CENTER, font_name="Courier New")
    sx18 += sw18 + sgap18

# Compatible Frameworks section
add_rect(sl, Inches(0.35), Inches(2.95), Inches(12.6), Inches(0.42), fill=TEAL)
add_textbox(sl, Inches(0.5), Inches(3.02), Inches(12), Inches(0.28),
            "Compatible Agent Frameworks", size=13, bold=True, color=WHITE)

frameworks = [
    ("Google ADK", "Agent Development Kit", NAVY,
     ["First-class A2A support", "Built-in A2A server/client", "Works with Gemini models"]),
    ("LangChain", "Popular LLM framework", BLUE2,
     ["A2A agent integration", "Tool and chain support", "Python and TypeScript"]),
    ("CrewAI", "Multi-agent orchestration", TEAL,
     ["Role-based agents", "Crew collaboration", "Task management"]),
    ("Semantic\nKernel", "Microsoft AI SDK", RGBColor(0x51, 0x2B, 0xD4),
     [".NET and Python", "Plugin architecture", "Enterprise focus"]),
    ("AutoGen", "Multi-agent conversations", GREEN,
     ["Conversational agents", "Code execution", "Group chat support"]),
    ("Custom\nPlatforms", "Build your own", MGRAY,
     ["Any HTTP server", "Any language/framework", "Just implement the spec"]),
]

fx18 = Inches(0.35)
fy18 = Inches(3.4)
fw18 = Inches(2.05)
fh18 = Inches(3.7)
fgap18 = Inches(0.08)

for fname18, ftype, fbg, ffeats in frameworks:
    add_rect(sl, fx18, fy18, fw18, fh18, fill=fbg)
    add_textbox(sl, fx18 + Inches(0.1), fy18 + Inches(0.1), fw18 - Inches(0.2), Inches(0.55),
                fname18, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, fx18 + Inches(0.1), fy18 + Inches(0.72), fw18 - Inches(0.2), Inches(0.3),
                ftype, size=10, italic=True, color=RGBColor(0xB0, 0xC4, 0xDE),
                align=PP_ALIGN.CENTER)
    add_rect(sl, fx18, fy18 + Inches(1.08), fw18, Inches(0.025), fill=GOLD)
    yf18 = fy18 + Inches(1.18)
    for ff in ffeats:
        add_textbox(sl, fx18 + Inches(0.12), yf18, fw18 - Inches(0.24), Inches(0.3),
                    "• " + ff, size=10.5, color=WHITE)
        yf18 += Inches(0.36)
    fx18 += fw18 + fgap18

add_notes(sl, "A2A has official SDKs for the five most popular languages: Python, JavaScript, Java, "
              "C#/.NET, and Go. All SDKs are open source and maintained on GitHub. The protocol "
              "works with major agent frameworks including Google ADK (first-class support), "
              "LangChain, CrewAI, Semantic Kernel, AutoGen, and any custom platform that implements "
              "the HTTP-based specification.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Key Takeaways
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Key Takeaways", "Why A2A matters for the future of AI agents")

takeaways = [
    ("1", "Open Standard, Not Proprietary",
     "A2A is governed by the Linux Foundation, ensuring vendor-neutral evolution. "
     "No single company controls the future of agent interoperability.",
     NAVY),
    ("2", "Agents as First-Class Citizens",
     "Unlike tool-wrapping, A2A treats agents as autonomous entities capable of "
     "multi-turn reasoning, negotiation, and long-running task management.",
     BLUE2),
    ("3", "Enterprise Ready from Day One",
     "OAuth 2.0, mTLS, OpenID Connect, distributed tracing, and monitoring are "
     "built in — not bolted on. Production deployments are well-supported.",
     TEAL),
    ("4", "Complements MCP, Doesn't Replace It",
     "MCP connects models to tools and data. A2A connects agents to agents. "
     "Together they form a complete AI integration stack.",
     GREEN),
    ("5", "Minimal New Concepts",
     "Built on HTTP, JSON-RPC, and SSE — technologies every developer knows. "
     "Low adoption cost with high interoperability payoff.",
     RGBColor(0x6A, 0x1B, 0x9A)),
    ("6", "The Foundation for Multi-Agent AI",
     "As AI systems grow more complex, standardized agent communication is "
     "essential. A2A is the protocol that makes collaborative AI systems possible.",
     RGBColor(0xE6, 0x5C, 0x00)),
]

tx19 = Inches(0.35)
ty19 = Inches(1.25)
tw19 = Inches(6.1)
th19 = Inches(1.55)
tgap19 = Inches(0.12)

for i, (tnum, ttitle, tdesc, tbg) in enumerate(takeaways):
    col_x = tx19 + (i % 2) * (tw19 + tgap19 + Inches(0.45))
    row_y = ty19 + (i // 2) * (th19 + tgap19)
    add_rect(sl, col_x, row_y, tw19, th19, fill=tbg)
    # Number badge
    add_rect(sl, col_x + Inches(0.12), row_y + Inches(0.15), Inches(0.42), Inches(0.42), fill=GOLD)
    add_textbox(sl, col_x + Inches(0.12), row_y + Inches(0.15), Inches(0.42), Inches(0.42),
                tnum, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(sl, col_x + Inches(0.65), row_y + Inches(0.15), tw19 - Inches(0.8), Inches(0.38),
                ttitle, size=13.5, bold=True, color=GOLD)
    add_textbox(sl, col_x + Inches(0.15), row_y + Inches(0.6), tw19 - Inches(0.3), Inches(0.85),
                tdesc, size=11.5, color=WHITE)

add_notes(sl, "Six key takeaways. A2A is an open, Linux Foundation-governed standard — not a "
              "proprietary protocol. It elevates agents from tools to first-class autonomous "
              "entities. It is enterprise-ready with full security and observability built in. "
              "It complements MCP in a layered AI stack. It has minimal adoption cost since it "
              "builds on familiar web standards. And it lays the foundation for the multi-agent "
              "AI systems that will power the next generation of AI applications.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Resources
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Resources", "Where to learn more and get started")

resources = [
    ("Official Protocol Site",  NAVY,  WHITE,
     "a2a-protocol.org",
     ["Full specification",
      "Topics & tutorials",
      "Reference implementations"]),
    ("Linux Foundation",        BLUE2, WHITE,
     "linuxfoundation.org",
     ["Governance & roadmap",
      "Community contribution",
      "Working group info"]),
    ("Official SDKs",           TEAL,  WHITE,
     "github.com/google-a2a",
     ["Python, JS, Java, C#, Go",
      "Open source (Apache 2.0)",
      "Examples & quickstarts"]),
    ("DeepLearning.AI Course",  GREEN, WHITE,
     "deeplearning.ai",
     ["A2A + ADK tutorial",
      "Hands-on coding labs",
      "Free short course"]),
]

rx20 = Inches(0.35)
ry20 = Inches(1.25)
rw20 = Inches(3.05)
rh20 = Inches(3.8)
rgap20 = Inches(0.12)

for rname, rbg, rfg, rurl, rpoints in resources:
    add_rect(sl, rx20, ry20, rw20, rh20, fill=rbg)
    add_textbox(sl, rx20 + Inches(0.12), ry20 + Inches(0.12), rw20 - Inches(0.24), Inches(0.45),
                rname, size=15, bold=True, color=GOLD)
    add_rect(sl, rx20, ry20 + Inches(0.62), rw20, Inches(0.025), fill=GOLD)
    add_textbox(sl, rx20 + Inches(0.12), ry20 + Inches(0.72), rw20 - Inches(0.24), Inches(0.35),
                rurl, size=12, bold=False, italic=True,
                color=RGBColor(0xB0, 0xC4, 0xDE), font_name="Courier New")
    yr20 = ry20 + Inches(1.18)
    for rp in rpoints:
        add_textbox(sl, rx20 + Inches(0.18), yr20, rw20 - Inches(0.3), Inches(0.32),
                    "• " + rp, size=12, color=WHITE)
        yr20 += Inches(0.38)
    rx20 += rw20 + rgap20

# Getting started box
add_rect(sl, Inches(0.35), Inches(5.2), Inches(12.6), Inches(1.95), fill=LGRAY, line=NAVY, line_w=Pt(1.5))
add_textbox(sl, Inches(0.55), Inches(5.3), Inches(5), Inches(0.35),
            "Quick Start (Python)", size=13, bold=True, color=NAVY)
qs_lines = [
    "pip install a2a-sdk",
    "# Build your first A2A agent:",
    "from a2a.server import A2AServer",
    "from a2a.types import AgentCard, AgentCapabilities",
]
yqs = Inches(5.7)
for ql in qs_lines:
    col_qs = RGBColor(0x28, 0x28, 0x28) if ql.startswith('#') else BLACK
    add_textbox(sl, Inches(0.55), yqs, Inches(5.8), Inches(0.28),
                ql, size=11, color=RGBColor(0x4A, 0x14, 0x8C) if ql.startswith('pip') else
                RGBColor(0x2E, 0x7D, 0x32) if ql.startswith('#') else NAVY,
                font_name="Courier New")
    yqs += Inches(0.27)

# QR-code placeholder
add_rect(sl, Inches(10.5), Inches(5.2), Inches(2.4), Inches(1.95), fill=WHITE, line=NAVY, line_w=Pt(1.5))
add_textbox(sl, Inches(10.6), Inches(5.35), Inches(2.2), Inches(0.35),
            "a2a-protocol.org", size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(10.6), Inches(5.75), Inches(2.2), Inches(1.0),
            "[ Scan or visit for\n  full documentation\n  and specification ]",
            size=10.5, color=MGRAY, align=PP_ALIGN.CENTER)

# Closing tagline
add_rect(sl, Inches(0.35), Inches(7.1), Inches(12.6), Inches(0.32), fill=NAVY)
add_textbox(sl, Inches(0.5), Inches(7.12), Inches(12.2), Inches(0.25),
            "The A2A Protocol — empowering AI agents to collaborate across every boundary.",
            size=12, bold=True, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

add_notes(sl, "Key resources for getting started with A2A. The official site at a2a-protocol.org "
              "has the full specification, topic guides, and tutorials. The Linux Foundation site "
              "has governance and community information. Official SDKs are on GitHub under "
              "google-a2a. DeepLearning.AI offers a free short course on building agents with "
              "A2A and the Google Agent Development Kit.")

# ── Save ──────────────────────────────────────────────────────────────────────
output_path = "/home/agenticai/dev/a2a/A2A_Protocol_Concepts.pptx"
prs.save(output_path)

import os
size_kb = os.path.getsize(output_path) / 1024
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
print(f"File size: {size_kb:.1f} KB  ({size_kb/1024:.2f} MB)")
